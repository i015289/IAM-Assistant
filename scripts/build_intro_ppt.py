"""
Build an 8-slide introduction PPT for the IAM Assistant.

Audience: business stakeholders / management.
Goal:     awareness and buy-in.
Tone:     story-driven, warm and clean.

Run with:
    conda run -n sapcli-env python scripts/build_intro_ppt.py
Output:
    docs/IAM_Assistant_Intro.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Theme
# ---------------------------------------------------------------------------

NAVY    = RGBColor(0x1E, 0x3A, 0x5F)
BLUE    = RGBColor(0x2E, 0x86, 0xC1)
ORANGE  = RGBColor(0xE6, 0x7E, 0x22)
GREEN   = RGBColor(0x27, 0xAE, 0x60)
DARK    = RGBColor(0x33, 0x33, 0x33)
MED     = RGBColor(0x66, 0x66, 0x66)
LIGHT   = RGBColor(0xEE, 0xEE, 0xEE)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BG      = RGBColor(0xF7, 0xF9, 0xFB)
WARM    = RGBColor(0xFD, 0xF6, 0xEC)   # warm off-white for story slides

W = Inches(13.333)
H = Inches(7.5)


# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------

def bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, l, t, w, h, fill=None, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.shadow.inherit = False
    if fill is not None:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.75)
    return s


def txt(slide, l, t, w, h, text, *, size=14, bold=False, color=DARK,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def bullets(slide, l, t, w, h, items, *, size=14, color=DARK,
            dot_color=BLUE, spacing=1.2):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = spacing
        p.space_after = Pt(5)
        dot = p.add_run()
        dot.text = "▸  "
        dot.font.name = "Calibri"
        dot.font.size = Pt(size)
        dot.font.bold = True
        dot.font.color.rgb = dot_color
        if isinstance(item, tuple):
            head, tail = item
            r1 = p.add_run()
            r1.text = head
            r1.font.name = "Calibri"
            r1.font.size = Pt(size)
            r1.font.bold = True
            r1.font.color.rgb = NAVY
            r2 = p.add_run()
            r2.text = "  " + tail
            r2.font.name = "Calibri"
            r2.font.size = Pt(size)
            r2.font.color.rgb = color
        else:
            r = p.add_run()
            r.text = item
            r.font.name = "Calibri"
            r.font.size = Pt(size)
            r.font.color.rgb = color
    return box


def header(slide, title, subtitle=None):
    rect(slide, Inches(0), Inches(0), Inches(0.18), H, fill=BLUE)
    txt(slide, Inches(0.5), Inches(0.22), Inches(12.3), Inches(0.65),
        title, size=28, bold=True, color=NAVY)
    if subtitle:
        txt(slide, Inches(0.5), Inches(0.88), Inches(12.3), Inches(0.42),
            subtitle, size=14, color=MED)
    rect(slide, Inches(0.5), Inches(1.4), Inches(12.5), Inches(0.04), fill=LIGHT)


def stat_box(slide, l, t, w, h, number, label, num_color=NAVY):
    rect(slide, l, t, w, h, fill=WHITE, line=LIGHT)
    txt(slide, l, t + Inches(0.2), w, Inches(0.85),
        number, size=44, bold=True, color=num_color, align=PP_ALIGN.CENTER)
    txt(slide, l, t + Inches(1.0), w, Inches(0.5),
        label, size=13, color=MED, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------

def slide_01_title(prs):
    """Title: IAM Assistant"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, NAVY)
    # Top accent bar
    rect(s, Inches(0), Inches(0), W, Inches(0.12), fill=BLUE)
    # Bottom accent bar
    rect(s, Inches(0), H - Inches(0.12), W, Inches(0.12), fill=ORANGE)
    # Main title
    txt(s, Inches(1.5), Inches(2.0), Inches(10.0), Inches(1.2),
        "IAM Assistant", size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Subtitle
    txt(s, Inches(1.5), Inches(3.3), Inches(10.0), Inches(0.7),
        "Smarter SAP IAM Analysis — powered by Claude AI",
        size=20, color=RGBColor(0xA8, 0xC8, 0xE8), align=PP_ALIGN.CENTER)
    # Tag line
    txt(s, Inches(1.5), Inches(4.3), Inches(10.0), Inches(0.5),
        "Treasury  ·  Cash Management  ·  Finance",
        size=14, color=RGBColor(0x78, 0xA0, 0xC0), align=PP_ALIGN.CENTER)


def slide_02_problem(prs):
    """The Problem: slow, manual, expert-dependent"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, WARM)
    header(s, "The Problem", "IAM analysis today is slow, manual, and expert-dependent")

    # Left column — pain points
    txt(s, Inches(0.5), Inches(1.6), Inches(6.5), Inches(0.45),
        "A typical investigation requires:", size=13, color=MED)
    bullets(s, Inches(0.5), Inches(2.1), Inches(6.2), Inches(3.8), [
        ("SAP developer access  ", "SE16, cross-table navigation, ABAP knowledge"),
        ("Manual cross-referencing  ", "auth objects, catalogs, BRTs, field values"),
        ("Deep domain expertise  ", "SoD rules, TRFCT semantics, IAM table structure"),
        ("Hours of effort  ", "for a question that should take minutes"),
    ], size=14, color=DARK)

    # Right column — big stat
    rect(s, Inches(7.8), Inches(1.6), Inches(4.8), Inches(4.5), fill=WHITE, line=LIGHT)
    txt(s, Inches(7.8), Inches(2.1), Inches(4.8), Inches(0.7),
        "Before IAM Assistant", size=13, bold=True, color=MED, align=PP_ALIGN.CENTER)
    txt(s, Inches(7.8), Inches(2.8), Inches(4.8), Inches(1.1),
        "Hours", size=52, bold=True, color=RGBColor(0xC0, 0x39, 0x2B), align=PP_ALIGN.CENTER)
    txt(s, Inches(7.8), Inches(3.9), Inches(4.8), Inches(0.5),
        "per investigation", size=14, color=MED, align=PP_ALIGN.CENTER)
    rect(s, Inches(9.0), Inches(4.6), Inches(2.4), Inches(0.04), fill=LIGHT)
    txt(s, Inches(7.8), Inches(4.8), Inches(4.8), Inches(0.8),
        "Requires an SAP consultant\nwith direct system access",
        size=12, color=MED, align=PP_ALIGN.CENTER)


def slide_03_what_it_is(prs):
    """What It Is"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, BG)
    header(s, "What Is IAM Assistant?", "A conversational AI that queries live SAP data in plain English")

    txt(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(0.5),
        "Ask questions in natural language. Get answers from live ER6 data — no SAP access, no SQL, no manual table navigation.",
        size=14, color=DARK)

    # Two mode cards
    for i, (title, icon, lines) in enumerate([
        ("Claude Code CLI", "⌨",  [
            "Specialist modes: /treasury-iam, /cash-iam, /fin-iam",
            "Persistent memo system — findings survive sessions",
            "Autonomous /execute agent for multi-step investigations",
            "Ideal for power users and analysts",
        ]),
        ("Web UI", "🌐", [
            "Browser-based chat — no CLI knowledge needed",
            "Streaming responses with live query progress",
            "Prompt library with ready-made investigation templates",
            "Shareable access across the team",
        ]),
    ]):
        lx = Inches(0.5) + i * Inches(6.5)
        rect(s, lx, Inches(2.3), Inches(6.0), Inches(4.0), fill=WHITE, line=LIGHT)
        txt(s, lx + Inches(0.25), Inches(2.5), Inches(5.5), Inches(0.6),
            f"{icon}  {title}", size=17, bold=True, color=NAVY)
        rect(s, lx + Inches(0.25), Inches(3.1), Inches(5.5), Inches(0.03), fill=LIGHT)
        bullets(s, lx + Inches(0.25), Inches(3.2), Inches(5.4), Inches(2.8),
                lines, size=13, color=DARK)


def slide_04_capabilities(prs):
    """Key Capabilities"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, BG)
    header(s, "Key Capabilities", "Built-in domain expertise across Treasury, Cash Management, and Finance")

    caps = [
        ("Treasury IAM",       NAVY,   "FOE / BOE / MOE SoD rules, TRFCT forbidden combinations, hedge request management, BRT audits"),
        ("Cash Management",    BLUE,   "Bank account lifecycle auth, four-eyes principle, submit/approve segregation, interest management"),
        ("Finance (AP/AR/GL)", GREEN,  "Accounts payable/receivable/general ledger auth objects, payment run SoD, journal entry controls"),
        ("Autonomous Agent",   ORANGE, "/execute chains multiple queries, adapts to findings, and delivers a structured report — no prompting needed"),
    ]

    for i, (cap, color, desc) in enumerate(caps):
        row = i // 2
        col = i % 2
        lx = Inches(0.5) + col * Inches(6.3)
        ty = Inches(1.7) + row * Inches(2.5)
        rect(s, lx, ty, Inches(6.0), Inches(2.2), fill=WHITE, line=LIGHT)
        rect(s, lx, ty, Inches(0.22), Inches(2.2), fill=color)
        txt(s, lx + Inches(0.4), ty + Inches(0.2), Inches(5.4), Inches(0.5),
            cap, size=16, bold=True, color=color)
        txt(s, lx + Inches(0.4), ty + Inches(0.75), Inches(5.3), Inches(1.3),
            desc, size=13, color=DARK)


def slide_05_example(prs):
    """Live Example: query → answer in 50–90 seconds"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, WARM)
    header(s, "In Practice", "A 5-query SoD investigation — from question to answer in under 90 seconds")

    # Question box
    rect(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(0.85), fill=NAVY)
    txt(s, Inches(0.8), Inches(1.72), Inches(11.8), Inches(0.65),
        "\"Does SAP_BR_CASH_MANAGER violate the four-eyes principle across submit and approve catalogs?\"",
        size=14, bold=True, color=WHITE)

    # Steps
    steps = [
        ("1", "Look up all catalogs assigned to the BRT",                       BLUE),
        ("2", "For each catalog, retrieve all apps and their auth objects",      BLUE),
        ("3", "Check F_CLM_BAOR ACTVT values — submit (01/06) vs approve (31)", BLUE),
        ("4", "Confirm four-eyes: submit and approve apps are in separate catalogs", GREEN),
        ("5", "Produce structured report: findings, verdict, recommendation",   GREEN),
    ]
    for i, (num, step, color) in enumerate(steps):
        ty = Inches(2.65) + i * Inches(0.72)
        rect(s, Inches(0.5), ty, Inches(0.52), Inches(0.55), fill=color)
        txt(s, Inches(0.5), ty, Inches(0.52), Inches(0.55),
            num, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, Inches(1.2), ty + Inches(0.06), Inches(10.5), Inches(0.45),
            step, size=13, color=DARK)

    # Timing stat
    rect(s, Inches(9.5), Inches(6.0), Inches(3.3), Inches(1.1), fill=NAVY)
    txt(s, Inches(9.5), Inches(6.05), Inches(3.3), Inches(0.55),
        "50 – 90 seconds", size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, Inches(9.5), Inches(6.6), Inches(3.3), Inches(0.4),
        "end-to-end investigation time", size=12, color=RGBColor(0xA8, 0xC8, 0xE8),
        align=PP_ALIGN.CENTER)


def slide_06_safety(prs):
    """Safety & Trust"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, BG)
    header(s, "Safe by Design", "Read-only access — no risk of accidental data modification")

    pillars = [
        ("Read-Only Access",     "🔒", GREEN,  "All queries run as the ANZEIGER display user. Write operations are impossible — even if requested."),
        ("SSL Encrypted",        "🔐", BLUE,   "All traffic between the assistant and ER6 is SSL-encrypted. Credentials never leave the system."),
        ("No Change Management", "✅", ORANGE, "The tool can be used freely during live operations — no change request, no system freeze required."),
        ("Full Audit Trail",     "📋", NAVY,   "Every ER6 query is logged with a timestamp to .session-log.md. Full replay and audit capability."),
    ]

    for i, (title, icon, color, desc) in enumerate(pillars):
        col = i % 2
        row = i // 2
        lx = Inches(0.5) + col * Inches(6.3)
        ty = Inches(1.7) + row * Inches(2.5)
        rect(s, lx, ty, Inches(6.0), Inches(2.1), fill=WHITE, line=LIGHT)
        txt(s, lx + Inches(0.3), ty + Inches(0.2), Inches(0.7), Inches(0.6),
            icon, size=28, align=PP_ALIGN.CENTER)
        txt(s, lx + Inches(1.1), ty + Inches(0.18), Inches(4.7), Inches(0.5),
            title, size=16, bold=True, color=color)
        txt(s, lx + Inches(1.1), ty + Inches(0.72), Inches(4.7), Inches(1.2),
            desc, size=13, color=DARK)


def slide_07_impact(prs):
    """Impact Summary"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, WARM)
    header(s, "Impact at a Glance", "What changes when your team uses IAM Assistant")

    # Stat row
    stat_box(s, Inches(0.5),  Inches(1.7), Inches(2.9), Inches(1.8), "~90 sec", "per investigation\n(was: hours)", num_color=GREEN)
    stat_box(s, Inches(3.6),  Inches(1.7), Inches(2.9), Inches(1.8), "3",       "domain specialist\nrole modes", num_color=BLUE)
    stat_box(s, Inches(6.7),  Inches(1.7), Inches(2.9), Inches(1.8), "100%",   "query test\npass rate (20 tests)", num_color=NAVY)
    stat_box(s, Inches(9.8),  Inches(1.7), Inches(2.9), Inches(1.8), "0",      "risk to live data\n(read-only ANZEIGER)", num_color=ORANGE)

    # Before / After
    for col, (title, color, items) in enumerate([
        ("Before", RGBColor(0xC0, 0x39, 0x2B), [
            "SAP developer required for every query",
            "Hours of manual table cross-referencing",
            "SoD rules in individual analysts' heads",
            "Findings lost between sessions",
        ]),
        ("After", GREEN, [
            "Any analyst asks in plain English",
            "50–90 seconds, fully automated",
            "Built-in rules: FOE/BOE, four-eyes, SoD",
            "Persistent memos — knowledge accumulates",
        ]),
    ]):
        lx = Inches(0.5) + col * Inches(6.5)
        rect(s, lx, Inches(3.8), Inches(6.0), Inches(0.45), fill=color)
        txt(s, lx, Inches(3.84), Inches(6.0), Inches(0.38),
            title, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        bullets(s, lx + Inches(0.2), Inches(4.35), Inches(5.6), Inches(2.5),
                items, size=13, color=DARK,
                dot_color=color)


def slide_08_next(prs):
    """What's Next / Q&A"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, NAVY)
    rect(s, Inches(0), Inches(0), W, Inches(0.12), fill=BLUE)
    rect(s, Inches(0), H - Inches(0.12), W, Inches(0.12), fill=ORANGE)

    txt(s, Inches(1.5), Inches(1.0), Inches(10.0), Inches(0.8),
        "What's Next", size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    items = [
        ("Status:  ",        "Live and running on ER6 today — Treasury, Cash Management, and Finance domains covered"),
        ("Access:  ",        "Web UI available via browser; Claude Code CLI for analysts who prefer terminal"),
        ("Domains:  ",       "Further expansion planned — AP/AR/GL coverage maturing, additional country variants"),
        ("Contribution:  ",  "Skill files are versioned — domain rules can be extended by the IAM team"),
    ]

    for i, (head, tail) in enumerate(items):
        ty = Inches(2.1) + i * Inches(0.95)
        rect(s, Inches(1.5), ty + Inches(0.08), Inches(0.06), Inches(0.5), fill=BLUE)
        txt(s, Inches(1.8), ty, Inches(2.0), Inches(0.6),
            head, size=15, bold=True, color=BLUE)
        txt(s, Inches(3.8), ty, Inches(8.5), Inches(0.6),
            tail, size=15, color=RGBColor(0xCC, 0xDD, 0xEE))

    rect(s, Inches(1.5), Inches(6.0), Inches(10.0), Inches(0.04), fill=BLUE)
    txt(s, Inches(1.5), Inches(6.15), Inches(10.0), Inches(0.5),
        "Questions?",
        size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def build():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    slide_01_title(prs)
    slide_02_problem(prs)
    slide_03_what_it_is(prs)
    slide_04_capabilities(prs)
    slide_05_example(prs)
    slide_06_safety(prs)
    slide_07_impact(prs)
    slide_08_next(prs)

    out = Path(__file__).parent.parent / "docs" / "IAM_Assistant_Intro.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    print(f"Saved → {out}")


if __name__ == "__main__":
    build()
