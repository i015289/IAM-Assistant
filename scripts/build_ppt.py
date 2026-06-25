"""
Build a PowerPoint presentation introducing the IAM Assistant.

Audience: mixed (technical + business).
Run with:
    conda run -n sapcli-env python scripts/build_ppt.py
Output:
    docs/IAM_Assistant_Overview.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu


# ---------- Theme (Catppuccin Latte/Mocha inspired, but PPT-friendly) ----------

NAVY = RGBColor(0x1E, 0x3A, 0x5F)        # primary dark
ACCENT = RGBColor(0x2E, 0x86, 0xC1)      # secondary blue
ACCENT2 = RGBColor(0xE6, 0x7E, 0x22)     # accent orange
GREEN = RGBColor(0x27, 0xAE, 0x60)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREY_DARK = RGBColor(0x33, 0x33, 0x33)
GREY_MED = RGBColor(0x66, 0x66, 0x66)
GREY_LIGHT = RGBColor(0xEE, 0xEE, 0xEE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_LIGHT = RGBColor(0xF7, 0xF9, 0xFB)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ---------- Helpers ----------

def set_slide_background(slide, color: RGBColor) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill=None, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    return shape


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    font_size: int = 14,
    bold: bool = False,
    color: RGBColor = GREY_DARK,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    font_name: str = "Calibri",
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor

    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(
    slide,
    left,
    top,
    width,
    height,
    items,
    *,
    font_size: int = 14,
    color: RGBColor = GREY_DARK,
    bullet_color: RGBColor = ACCENT,
    line_spacing: float = 1.15,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)

    for idx, item in enumerate(items):
        # item can be a string, or (head, tail) tuple for bold prefix + body
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(4)

        # Bullet glyph
        r0 = p.add_run()
        r0.text = "▸  "
        r0.font.name = "Calibri"
        r0.font.size = Pt(font_size)
        r0.font.bold = True
        r0.font.color.rgb = bullet_color

        if isinstance(item, tuple):
            head, tail = item
            r1 = p.add_run()
            r1.text = head
            r1.font.name = "Calibri"
            r1.font.size = Pt(font_size)
            r1.font.bold = True
            r1.font.color.rgb = NAVY

            r2 = p.add_run()
            r2.text = "  " + tail
            r2.font.name = "Calibri"
            r2.font.size = Pt(font_size)
            r2.font.color.rgb = color
        else:
            r1 = p.add_run()
            r1.text = item
            r1.font.name = "Calibri"
            r1.font.size = Pt(font_size)
            r1.font.color.rgb = color
    return box


def add_header_bar(slide, title: str, subtitle: str | None = None):
    # left accent strip
    add_rect(slide, Inches(0), Inches(0), Inches(0.18), SLIDE_H, fill=ACCENT)
    # title
    add_text(
        slide,
        Inches(0.5),
        Inches(0.25),
        Inches(12.3),
        Inches(0.6),
        title,
        font_size=28,
        bold=True,
        color=NAVY,
    )
    if subtitle:
        add_text(
            slide,
            Inches(0.5),
            Inches(0.85),
            Inches(12.3),
            Inches(0.4),
            subtitle,
            font_size=14,
            color=GREY_MED,
        )
    # underline
    add_rect(
        slide,
        Inches(0.5),
        Inches(1.30),
        Inches(1.2),
        Inches(0.05),
        fill=ACCENT2,
    )


def add_footer(slide, page_num: int, total: int):
    add_text(
        slide,
        Inches(0.5),
        Inches(7.10),
        Inches(8),
        Inches(0.3),
        "IAM Assistant  ·  Claude Code–powered IAM analysis on SAP ER6",
        font_size=10,
        color=GREY_MED,
    )
    add_text(
        slide,
        Inches(11.5),
        Inches(7.10),
        Inches(1.5),
        Inches(0.3),
        f"{page_num} / {total}",
        font_size=10,
        color=GREY_MED,
        align=PP_ALIGN.RIGHT,
    )


def add_card(slide, left, top, width, height, *, title, body, accent=ACCENT, body_size=12):
    add_rect(slide, left, top, width, height, fill=WHITE, line=GREY_LIGHT)
    add_rect(slide, left, top, width, Inches(0.07), fill=accent)
    add_text(
        slide,
        left + Inches(0.15),
        top + Inches(0.15),
        width - Inches(0.3),
        Inches(0.45),
        title,
        font_size=14,
        bold=True,
        color=NAVY,
    )
    add_text(
        slide,
        left + Inches(0.15),
        top + Inches(0.65),
        width - Inches(0.3),
        height - Inches(0.75),
        body,
        font_size=body_size,
        color=GREY_DARK,
    )


def add_table(slide, left, top, width, height, headers, rows,
              header_fill=NAVY, header_color=WHITE,
              body_size=11, header_size=12):
    n_cols = len(headers)
    n_rows = len(rows) + 1
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = tbl_shape.table

    # Header row
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        tf = cell.text_frame
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)
        tf.margin_top = Inches(0.05)
        tf.margin_bottom = Inches(0.05)
        tf.text = ""
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = h
        run.font.name = "Calibri"
        run.font.size = Pt(header_size)
        run.font.bold = True
        run.font.color.rgb = header_color

    # Body rows
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 == 1 else BG_LIGHT
            tf = cell.text_frame
            tf.margin_left = Inches(0.08)
            tf.margin_right = Inches(0.08)
            tf.margin_top = Inches(0.04)
            tf.margin_bottom = Inches(0.04)
            tf.word_wrap = True
            tf.text = ""
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(val)
            run.font.name = "Calibri"
            run.font.size = Pt(body_size)
            run.font.color.rgb = GREY_DARK
    return tbl


def add_code_block(slide, left, top, width, height, code: str,
                   bg=RGBColor(0x2D, 0x2D, 0x2D), fg=RGBColor(0xE6, 0xE6, 0xE6),
                   font_size=11):
    add_rect(slide, left, top, width, height, fill=bg)
    box = slide.shapes.add_textbox(
        left + Inches(0.1), top + Inches(0.08),
        width - Inches(0.2), height - Inches(0.16),
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)

    lines = code.splitlines() or [""]
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.05
        run = p.add_run()
        run.text = line if line else " "
        run.font.name = "Consolas"
        run.font.size = Pt(font_size)
        run.font.color.rgb = fg


# ---------- Slide builders ----------

def slide_title(prs, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_background(slide, NAVY)

    # Decorative bars
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.25), fill=ACCENT2)
    add_rect(slide, Inches(0), SLIDE_H - Inches(0.25), SLIDE_W, Inches(0.25), fill=ACCENT)

    # Logo-ish badge
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.0), Inches(1.4), Inches(2.4), Inches(0.55),
    )
    badge.shadow.inherit = False
    badge.fill.solid()
    badge.fill.fore_color.rgb = ACCENT2
    badge.line.fill.background()
    tf = badge.text_frame
    tf.text = ""
    tf.margin_left = Inches(0.15)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "SAP ER6  ·  IAM ANALYSIS"
    run.font.name = "Calibri"
    run.font.bold = True
    run.font.size = Pt(13)
    run.font.color.rgb = WHITE

    # Title
    add_text(
        slide,
        Inches(1.0),
        Inches(2.1),
        Inches(11.3),
        Inches(1.4),
        "IAM Assistant",
        font_size=64,
        bold=True,
        color=WHITE,
    )
    # Subtitle
    add_text(
        slide,
        Inches(1.0),
        Inches(3.5),
        Inches(11.3),
        Inches(0.8),
        "A Claude Code–powered assistant for analyzing",
        font_size=24,
        color=RGBColor(0xCF, 0xDB, 0xE8),
    )
    add_text(
        slide,
        Inches(1.0),
        Inches(3.95),
        Inches(11.3),
        Inches(0.8),
        "SAP IAM data on the ER6 system",
        font_size=24,
        color=RGBColor(0xCF, 0xDB, 0xE8),
    )

    # Tagline
    add_text(
        slide,
        Inches(1.0),
        Inches(5.2),
        Inches(11.3),
        Inches(0.5),
        "Natural language  →  multi-step ER6 queries  →  clean answers in minutes",
        font_size=16,
        color=ACCENT2,
        bold=True,
    )

    # Author/date
    add_text(
        slide,
        Inches(1.0),
        Inches(6.4),
        Inches(11.3),
        Inches(0.4),
        "June 2026",
        font_size=14,
        color=RGBColor(0xCF, 0xDB, 0xE8),
    )


def slide_agenda(prs, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_header_bar(slide, "Agenda", "What we will cover")

    items = [
        ("1.  The Problem", "Why IAM analysis on SAP is hard today"),
        ("2.  What is IAM Assistant", "Two modes — Claude Code CLI and Web UI"),
        ("3.  Architecture", "End-to-end flow from question to answer"),
        ("4.  Domain Skills", "Treasury, Cash Management, Finance specialists"),
        ("5.  Workflow Tools", "/goal, /execute, /memo and quality hooks"),
        ("6.  Live Examples", "SoD validation and catalog split scenarios"),
        ("7.  Benefits & Impact", "Time savings, safety, knowledge accumulation"),
        ("8.  Roadmap & Q&A", "What's next and open discussion"),
    ]

    box = slide.shapes.add_textbox(
        Inches(1.0), Inches(1.7), Inches(11.3), Inches(5.0)
    )
    tf = box.text_frame
    tf.word_wrap = True
    for idx, (head, tail) in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        p.line_spacing = 1.2

        r1 = p.add_run()
        r1.text = head
        r1.font.name = "Calibri"
        r1.font.size = Pt(20)
        r1.font.bold = True
        r1.font.color.rgb = NAVY

        r2 = p.add_run()
        r2.text = "    —    " + tail
        r2.font.name = "Calibri"
        r2.font.size = Pt(18)
        r2.font.color.rgb = GREY_DARK

    add_footer(slide, page_num, total)


def slide_problem(prs, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_header_bar(slide, "The Problem", "Analyzing SAP IAM today is slow, manual, and expert-bound")

    # Left card — pain points
    add_rect(slide, Inches(0.5), Inches(1.6), Inches(6.0), Inches(5.0),
             fill=WHITE, line=GREY_LIGHT)
    add_text(slide, Inches(0.7), Inches(1.75), Inches(5.6), Inches(0.5),
             "Today's reality", font_size=18, bold=True, color=RED)
    pain = [
        ("SAP expertise required.", "Direct ER6 access via SE16, plus deep table knowledge."),
        ("Cross-table joins by hand.", "APPAUI → APPAUV → BC_APP → BRTBUC..."),
        ("SoD rules in people's heads.", "FOE/BOE forbidden combos, TRFCT semantics, country variants."),
        ("Findings disappear.", "No persistent record between sessions or analysts."),
        ("Slow turnaround.", "A single SoD validation can take hours."),
    ]
    add_bullets(slide, Inches(0.7), Inches(2.3), Inches(5.6), Inches(4.2),
                pain, font_size=13, bullet_color=RED)

    # Right card — what's needed
    add_rect(slide, Inches(6.85), Inches(1.6), Inches(6.0), Inches(5.0),
             fill=WHITE, line=GREY_LIGHT)
    add_text(slide, Inches(7.05), Inches(1.75), Inches(5.6), Inches(0.5),
             "What we need", font_size=18, bold=True, color=GREEN)
    need = [
        ("Natural-language access.", "Ask in English; system handles SQL."),
        ("Encoded domain rules.", "SoD matrices, table quirks, query pitfalls."),
        ("Persistent memory.", "Findings survive across sessions."),
        ("Autonomous multi-step work.", "Chain queries; adapt to results."),
        ("Read-only safety.", "No risk to production IAM data."),
    ]
    add_bullets(slide, Inches(7.05), Inches(2.3), Inches(5.6), Inches(4.2),
                need, font_size=13, bullet_color=GREEN)

    add_footer(slide, page_num, total)


def slide_what_is(prs, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_header_bar(slide, "What is IAM Assistant?",
                   "An AI-powered analysis layer on top of SAP ER6")

    # Definition box
    add_rect(slide, Inches(0.5), Inches(1.55), Inches(12.3), Inches(1.3),
             fill=NAVY)
    add_text(
        slide, Inches(0.7), Inches(1.7), Inches(11.9), Inches(1.0),
        "An interactive assistant that translates natural-language IAM questions into "
        "correct multi-step ABAP Open SQL query chains against ER6, "
        "returning clean answers — and remembering them.",
        font_size=18, color=WHITE, bold=False,
    )

    # Two modes
    add_card(
        slide,
        Inches(0.5), Inches(3.05),
        Inches(6.0), Inches(3.7),
        title="Mode 1  ·  Claude Code CLI",
        body=(
            "Interactive command-line assistant powered by Claude Opus.\n\n"
            "•  Domain skills: /treasury-iam, /cash-iam, /fin-iam\n"
            "•  Wiki researcher: /iam-wiki (SAP Confluence via sap-wiki MCP)\n"
            "•  Goal-driven autonomous agent: /goal + /execute\n"
            "•  Persistent memo system: /memo\n"
            "•  Quality enforcement via three lifecycle hooks\n"
            "•  Direct MCP access to ER6 — no shell scripting\n\n"
            "Best for: deep multi-step investigations, audits, and split planning."
        ),
        accent=ACCENT,
        body_size=12,
    )

    add_card(
        slide,
        Inches(6.85), Inches(3.05),
        Inches(6.0), Inches(3.7),
        title="Mode 2  ·  Web UI",
        body=(
            "Browser-based chat UI built on FastAPI.\n\n"
            "•  Same Claude model, same MCP tools as the CLI\n"
            "•  ER6 + SAP Wiki access via MCPMultiClient (er6 + sap-wiki)\n"
            "•  Header shows ER6 and Wiki connection status independently\n"
            "•  OIDC authentication; per-user sessions\n"
            "•  Light / Dark theme (Catppuccin Latte / Mocha)\n"
            "•  Categorised prompt template library\n"
            "•  Streaming responses with markdown rendering\n\n"
            "Best for: shared team access, ad-hoc questions, and onboarding new analysts."
        ),
        accent=ACCENT2,
        body_size=12,
    )

    add_footer(slide, page_num, total)


def slide_architecture(prs, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_header_bar(slide, "Architecture",
                   "Natural language in, validated answers out")

    # Pipeline boxes
    boxes = [
        ("User", "Analyst asks a\nquestion in English", RGBColor(0x6C, 0x5C, 0xE7)),
        ("Claude (Opus)", "Selects skill;\nplans queries", ACCENT),
        ("Skills + Hooks", "Domain context\nand validation", ACCENT2),
        ("MCP Server", "ER6 tool layer\n(stdio JSON-RPC)", GREEN),
        ("SAP ER6", "ABAP Open SQL\n(read-only ANZEIGER)", NAVY),
    ]

    n = len(boxes)
    box_w_in = 2.2
    box_h_in = 1.5
    gap_in = 0.25
    total_w_in = box_w_in * n + gap_in * (n - 1)
    start_left_in = (13.333 - total_w_in) / 2
    top = Inches(2.0)
    box_w = Inches(box_w_in)
    box_h = Inches(box_h_in)
    gap = Inches(gap_in)

    for i, (title, body, color) in enumerate(boxes):
        left = Inches(start_left_in + (box_w_in + gap_in) * i)
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, box_w, box_h)
        rect.shadow.inherit = False
        rect.fill.solid()
        rect.fill.fore_color.rgb = color
        rect.line.fill.background()

        add_text(slide, left, top + Inches(0.15), box_w, Inches(0.45),
                 title, font_size=16, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)
        add_text(slide, left, top + Inches(0.65), box_w, Inches(0.85),
                 body, font_size=11, color=WHITE, align=PP_ALIGN.CENTER)

        # Arrow between boxes
        if i < n - 1:
            arrow_left = Inches(start_left_in + (box_w_in + gap_in) * i + box_w_in)
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                arrow_left + Inches(0.02),
                top + Inches(0.55),
                Inches(gap_in - 0.04),
                Inches(0.4),
            )
            arrow.shadow.inherit = False
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = GREY_MED
            arrow.line.fill.background()

    # Below — return path note
    add_text(slide, Inches(0.5), Inches(3.85), Inches(12.3), Inches(0.4),
             "Result rows + analysis stream back to the user with citations to the queries that produced them.",
             font_size=12, color=GREY_MED, align=PP_ALIGN.CENTER)

    # Components used
    add_text(slide, Inches(0.5), Inches(4.45), Inches(12.3), Inches(0.4),
             "Stack components", font_size=16, bold=True, color=NAVY)

    headers = ["Layer", "Component", "Notes"]
    rows = [
        ["Frontend (UI)", "FastAPI + Jinja2 + vanilla JS", "Light/Dark theme, streaming chat"],
        ["LLM", "Claude Opus (latest)", "Via Hyperspace proxy at localhost:6655"],
        ["Skills", "Markdown skill files", "Mirrored to .claude/skills, .claude/commands, skills/"],
        ["MCP Server", "er6_mcp_server.py", "stdio JSON-RPC; six ER6 tools"],
        ["SAP", "ER6 (ABAP)", "ANZEIGER user, read-only, SSL"],
        ["Hooks", "Bash hooks in .claude/", "validate-memo, sync-skills, log-query"],
    ]
    add_table(slide, Inches(0.5), Inches(4.85), Inches(12.3), Inches(2.0),
              headers, rows, body_size=11, header_size=12)

    add_footer(slide, page_num, total)


def slide_skills(prs, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_header_bar(slide, "Domain Skills",
                   "Three IAM specialists with full encoded domain knowledge")

    skills = [
        {
            "name": "/treasury-iam",
            "color": ACCENT,
            "scope": "Treasury (TRM)",
            "objects": "T_DEAL_PD / PF / DP / AG, T_TOE_HR",
            "rules": "FOE / BOE / MOE forbidden TRFCT × ACTVT matrices, hedge SoD",
            "covers": "30+ auth objects, BRT inventory, country variants",
        },
        {
            "name": "/cash-iam",
            "color": ACCENT2,
            "scope": "Cash Management",
            "objects": "F_CLM_BAM / BAOR / BKCR / BAI / BAIC / BAH2",
            "rules": "Four-eyes principle (submit ≠ approve)",
            "covers": "BAM2, BAA submit/approve, BAI catalogs",
        },
        {
            "name": "/fin-iam",
            "color": GREEN,
            "scope": "Finance — AP / AR / GL / BA",
            "objects": "F_BKPF_BUK / BLA / KOA, F_LFA1_BUK, F_KNA1_BUK",
            "rules": "Payment-run SoD, AR dunning, GL posting controls",
            "covers": "AP, AR, GL, BA packages",
        },
    ]

    card_w = Inches(4.0)
    card_h = Inches(4.7)
    gap = Inches(0.27)
    start = Inches(0.5)
    top = Inches(1.65)

    for i, s in enumerate(skills):
        left = start + (card_w + gap) * i
        add_rect(slide, left, top, card_w, card_h, fill=WHITE, line=GREY_LIGHT)
        add_rect(slide, left, top, card_w, Inches(0.6), fill=s["color"])
        add_text(slide, left + Inches(0.2), top + Inches(0.1),
                 card_w - Inches(0.4), Inches(0.45),
                 s["name"], font_size=20, bold=True, color=WHITE)

        # Body fields
        y = top + Inches(0.85)
        for label, value in [
            ("Scope", s["scope"]),
            ("Auth objects", s["objects"]),
            ("SoD rules", s["rules"]),
            ("Coverage", s["covers"]),
        ]:
            add_text(slide, left + Inches(0.2), y,
                     card_w - Inches(0.4), Inches(0.3),
                     label, font_size=11, bold=True, color=GREY_MED)
            add_text(slide, left + Inches(0.2), y + Inches(0.28),
                     card_w - Inches(0.4), Inches(0.7),
                     value, font_size=12, color=GREY_DARK)
            y += Inches(0.95)

    # Footer note
    add_text(slide, Inches(0.5), Inches(6.45), Inches(12.3), Inches(0.4),
             "Activated explicitly  (e.g. /treasury-iam)  or automatically by keyword match.",
             font_size=12, color=GREY_MED, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
             "Plus  /iam-wiki  — SAP Confluence wiki researcher (available in both CLI and Web UI; SAP internal network required).",
             font_size=11, color=GREY_MED, align=PP_ALIGN.CENTER)

    add_footer(slide, page_num, total)


def slide_workflow(prs, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_header_bar(slide, "Workflow Tools",
                   "Goal → Execute → Memo, with hooks enforcing quality")

    # Top — three workflow tools
    cards = [
        ("/goal", "Decompose a high-level objective into a step-by-step plan.",
         "Validate all apps in SAP_TC_FIN_TRM_COMMON for FOE/BOE compliance and propose a split.",
         ACCENT),
        ("/execute", "Run the plan autonomously, chain queries, adapt to findings.",
         "Produces a structured Goal / Summary / Findings / Violations / Recommendations report.",
         ACCENT2),
        ("/memo", "Save and resume investigations across sessions.",
         "Findings, decisions, work-in-progress, known-good baselines preserved as markdown.",
         GREEN),
    ]
    card_w = Inches(4.0)
    gap = Inches(0.27)
    start = Inches(0.5)
    top = Inches(1.6)

    for i, (name, desc, example, color) in enumerate(cards):
        left = start + (card_w + gap) * i
        add_rect(slide, left, top, card_w, Inches(2.6), fill=WHITE, line=GREY_LIGHT)
        add_rect(slide, left, top, Inches(0.1), Inches(2.6), fill=color)
        add_text(slide, left + Inches(0.25), top + Inches(0.15),
                 card_w - Inches(0.4), Inches(0.5),
                 name, font_size=20, bold=True, color=color)
        add_text(slide, left + Inches(0.25), top + Inches(0.7),
                 card_w - Inches(0.4), Inches(0.85),
                 desc, font_size=12, color=GREY_DARK)
        add_text(slide, left + Inches(0.25), top + Inches(1.55),
                 card_w - Inches(0.4), Inches(0.3),
                 "Example", font_size=10, bold=True, color=GREY_MED)
        add_text(slide, left + Inches(0.25), top + Inches(1.8),
                 card_w - Inches(0.4), Inches(0.75),
                 example, font_size=11, color=GREY_DARK)

    # Hooks panel
    add_text(slide, Inches(0.5), Inches(4.45), Inches(12.3), Inches(0.4),
             "Quality hooks  ·  enforced automatically",
             font_size=16, bold=True, color=NAVY)

    headers = ["Hook", "Trigger", "Effect"]
    rows = [
        ["validate-memo.sh", "PreToolUse / Write to .claude/memo/*.md",
         "Blocks writes if any of the 4 required sections is missing"],
        ["sync-skills.sh", "PostToolUse / Write to .claude/skills/*.md",
         "Auto-syncs skill to skills/ and .claude/commands/ — prevents drift"],
        ["log-query.sh", "PostToolUse / mcp__er6__query_sql",
         "Appends every ER6 query with timestamp to .session-log.md"],
    ]
    add_table(slide, Inches(0.5), Inches(4.85), Inches(12.3), Inches(1.85),
              headers, rows, body_size=11, header_size=12)

    add_footer(slide, page_num, total)


def slide_example(prs, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_header_bar(slide, "Example  ·  Treasury SoD Validation",
                   "From one English question to a validated answer in ~60 seconds")

    # Left — user prompt
    add_text(slide, Inches(0.5), Inches(1.55), Inches(6.0), Inches(0.4),
             "1.  User asks", font_size=14, bold=True, color=ACCENT)
    add_rect(slide, Inches(0.5), Inches(1.95), Inches(6.0), Inches(1.0),
             fill=WHITE, line=GREY_LIGHT)
    add_text(slide, Inches(0.65), Inches(2.05), Inches(5.7), Inches(0.85),
             "/treasury-iam\n"
             "For app FTRCAI02_B_TRAN, validate whether it is "
             "compliant with BOE SoD rules.",
             font_size=13, color=GREY_DARK)

    # Left — chain
    add_text(slide, Inches(0.5), Inches(3.1), Inches(6.0), Inches(0.4),
             "2.  Assistant runs (autonomously)", font_size=14, bold=True, color=ACCENT)

    chain = [
        "Look up app in APS_IAM_W_APP",
        "Find catalog via APS_IAM_W_BC_APP",
        "Read auth instances from APPAUI",
        "Pull TRFCT × ACTVT values from APPAUV",
        "Check against BOE forbidden matrix (D2 × 01/02/16/85/KU/VF)",
        "Cross-check BRT linkage in APS_IAM_W_BRTBUC",
    ]
    box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(6.0), Inches(3.4))
    tf = box.text_frame
    tf.word_wrap = True
    for idx, step in enumerate(chain):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.15
        p.space_after = Pt(4)
        r0 = p.add_run()
        r0.text = f"{idx + 1}.  "
        r0.font.name = "Calibri"
        r0.font.size = Pt(13)
        r0.font.bold = True
        r0.font.color.rgb = ACCENT
        r1 = p.add_run()
        r1.text = step
        r1.font.name = "Calibri"
        r1.font.size = Pt(13)
        r1.font.color.rgb = GREY_DARK

    # Right — verdict
    add_text(slide, Inches(6.85), Inches(1.55), Inches(6.0), Inches(0.4),
             "3.  Result  ·  ~60 seconds", font_size=14, bold=True, color=GREEN)

    add_rect(slide, Inches(6.85), Inches(1.95), Inches(6.0), Inches(4.95),
             fill=WHITE, line=GREY_LIGHT)
    add_rect(slide, Inches(6.85), Inches(1.95), Inches(6.0), Inches(0.08),
             fill=GREEN)
    add_text(slide, Inches(7.0), Inches(2.1), Inches(5.7), Inches(0.5),
             "Verdict  ·  COMPLIANT  ✓", font_size=18, bold=True, color=GREEN)

    verdict_lines = [
        ("Catalog", "SAP_FIN_BC_TRM_BACK_OFFICE_PC"),
        ("BRT", "SAP_BR_TREASURY_SPECIALIST_BOE"),
        ("Auth objects active", "T_DEAL_PD, T_DEAL_PF, T_DEAL_DP"),
        ("TRFCT values held", "D3 (Settlement only)"),
        ("ACTVT values held", "01, 02, 03, 16"),
        ("Forbidden combo D2 × 01", "Not present"),
        ("Forbidden combo D2 × 16", "Not present"),
        ("Cross-BRT exposure", "None  (BOE only)"),
    ]
    y = Inches(2.7)
    for label, value in verdict_lines:
        add_text(slide, Inches(7.0), y, Inches(2.7), Inches(0.3),
                 label, font_size=11, bold=True, color=GREY_MED)
        add_text(slide, Inches(9.7), y, Inches(3.1), Inches(0.3),
                 value, font_size=11, color=GREY_DARK)
        y += Inches(0.5)

    add_text(slide, Inches(7.0), Inches(6.5), Inches(5.7), Inches(0.3),
             "All findings cite the exact ER6 queries that produced them (auditable).",
             font_size=10, color=GREY_MED)

    add_footer(slide, page_num, total)


def slide_benefits(prs, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_header_bar(slide, "Benefits", "Concrete, measured impact")

    # Big metrics row
    metrics = [
        ("50–90 s", "Typical 5–8 query investigation\nround-trip time", ACCENT),
        ("100 %", "Pass rate on validated\nquery patterns (20 tests)", GREEN),
        ("0", "Risk to live data\n(read-only ANZEIGER user)", ACCENT2),
        ("3", "Quality hooks running\nautomatically every session", NAVY),
    ]
    card_w = Inches(2.95)
    gap = Inches(0.2)
    start = Inches(0.55)
    top = Inches(1.6)

    for i, (big, small, color) in enumerate(metrics):
        left = start + (card_w + gap) * i
        add_rect(slide, left, top, card_w, Inches(1.7), fill=WHITE, line=GREY_LIGHT)
        add_rect(slide, left, top, Inches(0.08), Inches(1.7), fill=color)
        add_text(slide, left, top + Inches(0.15), card_w, Inches(0.85),
                 big, font_size=44, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(slide, left, top + Inches(1.05), card_w, Inches(0.65),
                 small, font_size=12, color=GREY_DARK, align=PP_ALIGN.CENTER)

    # Two-column benefits
    add_text(slide, Inches(0.55), Inches(3.55), Inches(12.3), Inches(0.4),
             "Why teams adopt it", font_size=16, bold=True, color=NAVY)

    left_items = [
        ("Speed.", "Hours of SE16 navigation collapse to a single conversation."),
        ("No SAP expertise required.", "Analysts ask in English; the assistant handles SQL and joins."),
        ("Encoded SoD knowledge.", "FOE/BOE/MOE matrices, TRFCT semantics, country variants — all built in."),
        ("Auditability.", "Every query is logged with timestamp; answers cite their queries."),
        ("Read-only by construction.", "ANZEIGER user, SSL — safe to use during operations."),
    ]
    add_bullets(slide, Inches(0.55), Inches(4.0), Inches(6.1), Inches(2.9),
                left_items, font_size=12, bullet_color=ACCENT)

    right_items = [
        ("Persistent memory.", "Multi-day investigations resume cleanly; no lost context."),
        ("Reusable skills.", "New analysts get full specialist context with /treasury-iam etc."),
        ("Autonomous execution.", "/execute chains queries, adapts to findings, returns structured reports."),
        ("Living documentation.", "Skill files and CLAUDE.md accumulate institutional knowledge."),
        ("Two access paths.", "Power users use the CLI; the team uses the Web UI."),
    ]
    add_bullets(slide, Inches(6.85), Inches(4.0), Inches(6.1), Inches(2.9),
                right_items, font_size=12, bullet_color=GREEN)

    add_footer(slide, page_num, total)


def slide_use_cases(prs, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_header_bar(slide, "Use Cases",
                   "What you can ask — by domain")

    headers = ["Domain", "Question you can ask", "What the assistant does"]
    rows = [
        ["Treasury — FOE/BOE",
         "Validate FTRCAI02_B_TRAN against BOE SoD rules.",
         "Reads APP / catalog / APPAUI / APPAUV; checks D3 forbidden combos; reports verdict."],
        ["Treasury — split",
         "Catalog SAP_TC_FIN_TRM_COMMON has both D2+01 and D3+01. Propose a split.",
         "Identifies offending apps; groups by office; produces FOE / BOE catalog plan."],
        ["Treasury — hedge",
         "Validate T_TOE_HR for SAP_FIN_BC_TRM_HM_HR_FOE_PC.",
         "Pulls HREQ_CAT × ACTVT; checks MOE-forbidden release/reverse on A and D."],
        ["Cash — four-eyes",
         "Does SAP_BR_CASH_MANAGER violate four-eyes between submit and approve?",
         "Maps BRT → catalogs; cross-checks F_CLM_BAOR ACTVT 31 vs submit catalog."],
        ["Cash — health check",
         "Run a full health check on F9017_TRAN.",
         "Lists active auth objects, ACTVT values, catalog assignments, BRT coverage."],
        ["Finance — AP",
         "Trace AP payment flow F0770 → F0771 → F0673A and confirm SoD.",
         "Walks F_BKPF_* and F_LFA1_* across propose / revise / approve."],
        ["Finance — GL",
         "GL journal entry: post / park / verify — confirm posting controls.",
         "Maps F_BKPF_BUK / BLA / KOA across the three apps."],
        ["Cross-cutting",
         "Find all Treasury apps holding D2 and D3 on T_DEAL_PD.",
         "APPAUV scan filtered by TRFCT, joined back to APP and catalog."],
    ]
    add_table(slide, Inches(0.5), Inches(1.6), Inches(12.3), Inches(5.3),
              headers, rows, body_size=10, header_size=11)

    add_footer(slide, page_num, total)


def slide_quickstart(prs, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_header_bar(slide, "Quick Start", "Three commands to get going")

    # Step 1
    add_text(slide, Inches(0.5), Inches(1.55), Inches(6.0), Inches(0.4),
             "1.  Install", font_size=16, bold=True, color=ACCENT)
    add_code_block(slide, Inches(0.5), Inches(2.0), Inches(6.0), Inches(0.9),
                   "$ ./install.sh           # macOS / Linux\n"
                   "$ install.bat            # Windows")

    # Step 2
    add_text(slide, Inches(0.5), Inches(3.0), Inches(6.0), Inches(0.4),
             "2.  Configure  (.env)", font_size=16, bold=True, color=ACCENT)
    add_code_block(slide, Inches(0.5), Inches(3.45), Inches(6.0), Inches(1.5),
                   "ANTHROPIC_API_KEY=<hyperspace-key>\n"
                   "ANTHROPIC_BASE_URL=http://localhost:6655\n"
                   "OIDC_CLIENT_ID=your-client-id\n"
                   "SESSION_SECRET=<long-random-string>")

    # Step 3
    add_text(slide, Inches(0.5), Inches(5.05), Inches(6.0), Inches(0.4),
             "3.  Run", font_size=16, bold=True, color=ACCENT)
    add_code_block(slide, Inches(0.5), Inches(5.5), Inches(6.0), Inches(1.4),
                   "# Web UI\n"
                   "$ conda run -n sapcli-env \\\n"
                   "    uvicorn app.main:app --reload\n"
                   "# Then open http://localhost:8080")

    # Right column — first questions to try
    add_text(slide, Inches(6.85), Inches(1.55), Inches(6.0), Inches(0.4),
             "First questions to try", font_size=16, bold=True, color=GREEN)
    add_rect(slide, Inches(6.85), Inches(2.0), Inches(6.0), Inches(4.9),
             fill=WHITE, line=GREY_LIGHT)

    examples = [
        "What can I ask?",
        "Glossary: what is a BRT, BC, App?",
        "What is SoD and why does it matter?",
        "/treasury-iam  Validate FTRCAI02_B_TRAN for BOE.",
        "/cash-iam  Check four-eyes for SAP_BR_CASH_MANAGER.",
        "/fin-iam  Trace the AP payment flow.",
        "/goal  Audit FOE/BOE for all apps in CLOUD_FI_TR_IAM.",
        "/execute  (after a /goal plan is presented)",
        "/memo save trm-foe-boe-validation",
    ]
    box = slide.shapes.add_textbox(Inches(7.05), Inches(2.15), Inches(5.7), Inches(4.6))
    tf = box.text_frame
    tf.word_wrap = True
    for idx, ex in enumerate(examples):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.25
        p.space_after = Pt(4)
        r0 = p.add_run()
        r0.text = "›  "
        r0.font.name = "Calibri"
        r0.font.size = Pt(12)
        r0.font.bold = True
        r0.font.color.rgb = GREEN
        r1 = p.add_run()
        r1.text = ex
        r1.font.name = "Consolas"
        r1.font.size = Pt(12)
        r1.font.color.rgb = GREY_DARK

    add_footer(slide, page_num, total)


def slide_roadmap(prs, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_header_bar(slide, "Roadmap & Open Items",
                   "Where the project goes next")

    # Three-column roadmap
    cols = [
        ("Now", ACCENT, [
            "Three domain skills live: Treasury, Cash, Finance (AP/AR/GL/BA).",
            "Web UI with OIDC auth, theme toggle, prompt template library.",
            "Goal/Execute/Memo workflow with three quality hooks.",
            "20-test query pattern validation (100% pass).",
        ]),
        ("Next", ACCENT2, [
            "More domains: Procurement, HCM, Group Reporting authorizations.",
            "Saved investigation library shared across the team.",
            "Inline charts: catalog coverage, BRT footprint heatmaps.",
            "Diff mode: pre / post-migration auth comparison reports.",
        ]),
        ("Later", GREEN, [
            "Multi-system support (DEV / QAS / PRD switching).",
            "Scheduled audits with email/Slack delivery.",
            "Self-service onboarding wizard for new analysts.",
            "Optional write-mode (with approval gating) for remediation.",
        ]),
    ]

    card_w = Inches(4.0)
    gap = Inches(0.27)
    start = Inches(0.5)
    top = Inches(1.6)
    height = Inches(5.0)

    for i, (label, color, items) in enumerate(cols):
        left = start + (card_w + gap) * i
        add_rect(slide, left, top, card_w, height, fill=WHITE, line=GREY_LIGHT)
        add_rect(slide, left, top, card_w, Inches(0.55), fill=color)
        add_text(slide, left + Inches(0.2), top + Inches(0.08),
                 card_w - Inches(0.4), Inches(0.45),
                 label, font_size=18, bold=True, color=WHITE)

        bullets_box = slide.shapes.add_textbox(
            left + Inches(0.2), top + Inches(0.7),
            card_w - Inches(0.4), height - Inches(0.85),
        )
        tf = bullets_box.text_frame
        tf.word_wrap = True
        for idx, item in enumerate(items):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.line_spacing = 1.2
            p.space_after = Pt(8)
            r0 = p.add_run()
            r0.text = "•  "
            r0.font.name = "Calibri"
            r0.font.bold = True
            r0.font.size = Pt(13)
            r0.font.color.rgb = color
            r1 = p.add_run()
            r1.text = item
            r1.font.name = "Calibri"
            r1.font.size = Pt(12)
            r1.font.color.rgb = GREY_DARK

    add_footer(slide, page_num, total)


def slide_thanks(prs, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, NAVY)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.25), fill=ACCENT2)
    add_rect(slide, Inches(0), SLIDE_H - Inches(0.25), SLIDE_W, Inches(0.25), fill=ACCENT)

    add_text(slide, Inches(1), Inches(2.5), Inches(11.3), Inches(1.4),
             "Thank you", font_size=72, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1), Inches(4.0), Inches(11.3), Inches(0.7),
             "Questions, demos, and prompt experiments welcome.",
             font_size=22, color=RGBColor(0xCF, 0xDB, 0xE8),
             align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1), Inches(5.0), Inches(11.3), Inches(0.5),
             "IAM Assistant  ·  built with Claude Code on SAP ER6",
             font_size=14, color=ACCENT2,
             align=PP_ALIGN.CENTER)


def slide_section(prs, page_num, total, label: str, title: str, color=ACCENT):
    """Section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, NAVY)
    add_rect(slide, Inches(0), Inches(3.4), SLIDE_W, Inches(0.08), fill=color)
    add_text(slide, Inches(1), Inches(2.4), Inches(11.3), Inches(0.6),
             label, font_size=18, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1), Inches(3.6), Inches(11.3), Inches(1.4),
             title, font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_footer(slide, page_num, total)


def slide_glossary(prs, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_header_bar(slide, "Glossary",
                   "Core IAM terms used throughout this deck")

    headers = ["Term", "Acronym", "Definition", "Example"]
    rows = [
        ["Business Role Template", "BRT",
         "Top-level role assigned to users; collection of business catalogs.",
         "SAP_BR_TREASURY_SPECIALIST_FOE"],
        ["Business Catalog", "BC",
         "Group of Fiori apps with shared restriction types and access rules.",
         "SAP_FIN_BC_TRM_BACK_OFFICE_PC"],
        ["IAM App", "App",
         "Atomic Fiori application unit with auth-object instances and values.",
         "FTRCAI02_B_TRAN"],
        ["Authorization Object", "Auth Obj",
         "ABAP entity controlling access; carries fields and activity values.",
         "T_DEAL_PD"],
        ["Activity", "ACTVT",
         "Standard SAP activity code (01=Create, 02=Change, 03=Display, …).",
         "01, 02, 16, 85"],
        ["Restriction Type", "RT",
         "Reusable scope-restriction grouping fields like company code.",
         "BUKRS_GLRLDNR"],
        ["Segregation of Duties", "SoD",
         "Control: same person cannot both perform and approve a sensitive action.",
         "FOE cannot Settle (D3 × 01)"],
        ["Front / Back / Middle Office", "FOE / BOE / MOE",
         "Treasury organisational split — trading vs settlement vs risk.",
         "FOE = traders; BOE = settlers"],
        ["Transaction Function Type", "TRFCT",
         "Treasury deal lifecycle stage on T_DEAL_*.",
         "D1=Order, D2=Contract, D3=Settle"],
        ["Hedge Request Category", "HREQ_CAT",
         "Hedge management lifecycle category on T_TOE_HR.",
         "A=Designation, D=Dedesignation"],
        ["Model Context Protocol", "MCP",
         "Open standard for connecting LLM agents to external tool servers.",
         "er6_mcp_server.py"],
        ["ER6 ANZEIGER", "—",
         "Read-only ABAP user used by IAM Assistant; SSL-encrypted.",
         "Display privileges only"],
    ]
    add_table(slide, Inches(0.4), Inches(1.55), Inches(12.5), Inches(5.45),
              headers, rows, body_size=10, header_size=11)
    add_footer(slide, page_num, total)


def slide_data_dict_overview(prs, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_header_bar(slide, "Data Dictionary  ·  Overview",
                   "The ER6 tables that drive every IAM analysis")

    headers = ["Table", "Domain", "Purpose", "Key Columns"]
    rows = [
        ["TDEVC", "Foundation", "ABAP packages",
         "DEVCLASS, PACKTYPE"],
        ["TADIR", "Foundation", "ABAP object registry — find SIA6 app objects",
         "OBJ_NAME, OBJECT, PGMID"],
        ["USOBT / USOBX", "Auth defaults", "T-code → required auth objects mapping",
         "NAME (TCODE), OBJECT, TYPE"],
        ["APS_IAM_W_APP", "App registry", "Master list of IAM apps",
         "APP_ID, APP_TYPE, TCODE, READ_ONLY"],
        ["APS_IAM_W_APPT", "App registry", "App display texts (multi-language)",
         "APP_ID, LANGU, TEXT"],
        ["APS_IAM_W_APPAUI", "App authz", "Auth object instances per app (UUID-keyed)",
         "APP_ID, UUID, AUTH_OBJECT, INACTIVE"],
        ["APS_IAM_W_APPAUV", "App authz", "Field-level values (TRFCT, ACTVT, …)",
         "APP_ID, UUID, PARENT_ID, FIELD, LOW/HIGH"],
        ["APS_IAM_W_APPAUO", "App authz", "Auth object exclusions (outbound)",
         "APP_ID, AUTH_OBJECT, STATUS"],
        ["APS_IAM_W_BC_APP", "Catalog", "Business catalog ↔ app assignments",
         "BU_CATALOG_ID, APP_ID, IS_FOLDER"],
        ["APS_IAM_W_BUC", "Catalog", "Business Catalog master  (⚠ use CDS  I_APS_BUSINESS_CATALOG)",
         "BU_CATALOG_ID, AGR_NAME"],
        ["APS_IAM_W_BRTBUC", "Role", "BRT ↔ Business Catalog assignments",
         "BRT_ID, BU_CATALOG_ID"],
        ["APS_IAM_W_BRT", "Role", "Business Role Template master",
         "BRT_ID, SCOPE_DEPENDENT, FIORI_SPACE_ID"],
        ["APS_IAM_W_RT / RR", "Restriction", "Restriction-type and restriction-rule masters",
         "RTYPE_ID, AUTH_OBJECT"],
        ["SUI_TM_MM_APP / CAT", "Launchpad", "Fiori Launchpad app↔catalog mapping",
         "APP_ID (GUID for TRM), CAT_ID"],
    ]
    add_table(slide, Inches(0.4), Inches(1.55), Inches(12.5), Inches(5.45),
              headers, rows, body_size=10, header_size=11)
    add_footer(slide, page_num, total)


def slide_table_relationships(prs, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_header_bar(slide, "Table Relationships",
                   "How a question flows through the schema")

    # ER-style boxes with arrows
    def er_box(left_in, top_in, w_in, h_in, title, fields, color=ACCENT):
        left = Inches(left_in)
        top = Inches(top_in)
        w = Inches(w_in)
        h = Inches(h_in)
        add_rect(slide, left, top, w, h, fill=WHITE, line=GREY_LIGHT)
        add_rect(slide, left, top, w, Inches(0.42), fill=color)
        add_text(slide, left + Inches(0.1), top + Inches(0.06),
                 w - Inches(0.2), Inches(0.32),
                 title, font_size=12, bold=True, color=WHITE)
        add_text(slide, left + Inches(0.12), top + Inches(0.5),
                 w - Inches(0.2), h - Inches(0.55),
                 fields, font_size=10, color=GREY_DARK)

    # Layer 1: BRT
    er_box(0.4, 1.65, 3.0, 0.9, "APS_IAM_W_BRT  ·  BRT master",
           "BRT_ID, FIORI_SPACE_ID, SCOPE_DEPENDENT", color=NAVY)

    # Layer 2: BRTBUC
    er_box(0.4, 2.85, 3.0, 0.9, "APS_IAM_W_BRTBUC  ·  BRT→BC",
           "BRT_ID, BU_CATALOG_ID", color=ACCENT)

    # Layer 3: BC_APP
    er_box(0.4, 4.05, 3.0, 0.9, "APS_IAM_W_BC_APP  ·  BC→App",
           "BU_CATALOG_ID, APP_ID, IS_FOLDER", color=ACCENT)

    # Layer 4: APP
    er_box(4.7, 4.05, 3.0, 0.9, "APS_IAM_W_APP  ·  App",
           "APP_ID, TCODE, APP_TYPE, READ_ONLY", color=ACCENT2)

    # Layer 5: APPAUI
    er_box(4.7, 5.25, 3.0, 0.9, "APS_IAM_W_APPAUI  ·  Instances",
           "APP_ID, UUID, AUTH_OBJECT, INACTIVE", color=ACCENT2)

    # Layer 6: APPAUV
    er_box(9.0, 5.25, 4.0, 0.9, "APS_IAM_W_APPAUV  ·  Field values",
           "APP_ID, UUID, PARENT_ID, FIELD, LOW_VALUE, HIGH_VALUE", color=GREEN)

    # Layer 7: AUTH OBJECT REGISTRY (right)
    er_box(9.0, 4.05, 4.0, 0.9, "APS_IAM_W_AU  ·  Auth obj registry",
           "AUTH_OBJECT_ID, AUTH_OBJECT, ACTIVITY_FIELD", color=GREEN)

    # Connecting arrows (visual only)
    def arrow(x1, y1, x2, y2):
        connector = slide.shapes.add_connector(
            1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        connector.line.color.rgb = GREY_MED
        connector.line.width = Pt(1.5)

    arrow(1.9, 2.55, 1.9, 2.85)   # BRT → BRTBUC
    arrow(1.9, 3.75, 1.9, 4.05)   # BRTBUC → BC_APP
    arrow(3.4, 4.5, 4.7, 4.5)     # BC_APP → APP
    arrow(6.2, 4.95, 6.2, 5.25)   # APP → APPAUI
    arrow(7.7, 5.7, 9.0, 5.7)     # APPAUI → APPAUV
    arrow(9.0, 4.5, 7.7, 4.5)     # AU registry ← APP

    # Legend / typical chain
    add_text(slide, Inches(0.4), Inches(6.4), Inches(12.5), Inches(0.4),
             "Typical SoD validation chain  ·  BRT  →  BRTBUC  →  BC_APP  →  APP  →  APPAUI  →  APPAUV  →  rule check",
             font_size=12, bold=True, color=NAVY)
    add_text(slide, Inches(0.4), Inches(6.75), Inches(12.5), Inches(0.3),
             "PARENT_ID on APPAUV links each value back to its APPAUI instance via UUID.",
             font_size=10, color=GREY_MED)
    add_footer(slide, page_num, total)


def slide_cds_view_layer(prs, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_header_bar(slide, "CDS View Layer  ·  Recommended over raw joins",
                   "1263 CDS views across 75 sub-packages under SR_APS_IAM_*  —  the supported semantic layer")

    # High-value views table
    add_text(slide, Inches(0.4), Inches(1.55), Inches(12.5), Inches(0.4),
             "High-value views for multi-hop questions",
             font_size=15, bold=True, color=NAVY)

    headers = ["CDS view", "Pre-joins / replaces", "Note"]
    rows = [
        ["I_APS_BUSINESS_CATALOG",
         "APS_IAM_W_BUC + 12 associations  (_App, _Successor, _RestrictionType, _BRT, …)",
         "Search-enabled. Always prefer over raw W_BUC."],
        ["APS_IAM_AUTH_FIELD_VAL",
         "APPAUI ⋈ APPAUV ⋈ TACTT ⋈ RT_AO  (App → AuthObject → Field → Value)",
         "ActivityValue is pre-translated. No scope filter — works for SIA1 + SIA6."],
        ["APS_IAM_INFO_BRT_BC_BASIC",
         "BRTBUC ⋈ BRT ⋈ BRTT  (BRT → BC with template name + component)",
         "Hidden filter: scope_state='3' or non-scope-dependent. Country variants may be missing — fall back to raw BRTBUC."],
    ]
    add_table(slide, Inches(0.4), Inches(2.0), Inches(12.5), Inches(1.95),
              headers, rows, body_size=10, header_size=11)

    # Routing rule
    add_text(slide, Inches(0.4), Inches(4.15), Inches(12.5), Inches(0.4),
             "⚠ Routing rule  ·  SIA6 (migrated) vs SIA1 (legacy) catalogs",
             font_size=15, bold=True, color=ACCENT2)

    routing = [
        ("Catalog migrated to IAM-Apps (SIA6).",
         "Use APS_IAM_BRT_APP for one-hop BRT → App."),
        ("Catalog still on SIA1  (Cash Management, Treasury today).",
         "BRT → BC via APS_IAM_INFO_BRT_BC_BASIC, then BC → App via raw APS_IAM_W_BC_APP. Auth view (APS_IAM_AUTH_FIELD_VAL) works regardless."),
    ]
    add_bullets(slide, Inches(0.4), Inches(4.6), Inches(12.5), Inches(1.1),
                routing, font_size=11, bullet_color=ACCENT2)

    # VDM naming convention card
    add_text(slide, Inches(0.4), Inches(5.8), Inches(12.5), Inches(0.4),
             "VDM naming convention",
             font_size=15, bold=True, color=NAVY)

    naming = [
        ("I_  Interface (reuse entry point).", "I_APS_BUSINESS_CATALOG, I_APS_IAM_APP_CORE"),
        ("C_  Consumption (OData / Fiori).", "C_APS_IAM_BR"),
        ("R_ / P_ / D_  Restricted / Private / Draft.", "R_APS_IAM_APP_TIL, P_APS_IAM_APP, D_APS_IAM_BUSR_RAP_*"),
        ("APS_IAM_*_DDL, *_BASIC.", "Pre-VDM legacy — still widely used."),
    ]
    add_bullets(slide, Inches(0.4), Inches(6.2), Inches(12.5), Inches(0.85),
                naming, font_size=10, bullet_color=ACCENT)
    add_footer(slide, page_num, total)


def slide_mcp_tools(prs, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_header_bar(slide, "MCP Tool Layer",
                   "ER6 data tools + SAP Wiki research tools — all over stdio")

    headers = ["Tool", "Purpose", "Typical use", "Mean latency"]
    rows = [
        ["mcp__er6__query_sql",
         "Run ABAP Open SQL SELECT statements",
         "Auth value lookups, BRT mapping, joins by hand",
         "10.1 s"],
        ["mcp__er6__read_table_def",
         "Read DDIC table / structure DDL",
         "Verify column names before writing SQL",
         "8.9 s"],
        ["mcp__er6__read_cds_view",
         "Read CDS view (DDLS) source",
         "Inspect view logic, parameters",
         "~9 s"],
        ["mcp__er6__read_class",
         "Read ABAP class source",
         "Trace BAdI / handler implementation",
         "~9 s"],
        ["mcp__er6__read_program",
         "Read program / report source",
         "Inspect helper reports, generated programs",
         "~9 s"],
        ["mcp__er6__list_package",
         "List all objects in a package",
         "Discover SIA6 apps in CLOUD_FI_TR_IAM",
         "17.3 s"],
    ]
    add_table(slide, Inches(0.4), Inches(1.55), Inches(12.5), Inches(2.6),
              headers, rows, body_size=11, header_size=12)

    # sap-wiki tools
    add_text(slide, Inches(0.4), Inches(4.3), Inches(12.5), Inches(0.35),
             "sap-wiki MCP  ·  SAP Confluence research (available in CLI and Web UI — SAP internal network required)",
             font_size=13, bold=True, color=NAVY)
    wiki_headers = ["Tool", "Purpose"]
    wiki_rows = [
        ["mcp__sap-wiki__general_search", "Keyword search across SAP Confluence"],
        ["mcp__sap-wiki__cql_search",     "Advanced CQL query search"],
        ["mcp__sap-wiki__wiki_content",   "Fetch full page content by URL"],
        ["mcp__sap-wiki__cql_examples",   "Get CQL syntax examples and reference"],
    ]
    add_table(slide, Inches(0.4), Inches(4.7), Inches(12.5), Inches(1.55),
              wiki_headers, wiki_rows, body_size=11, header_size=12)

    add_footer(slide, page_num, total)



def slide_skill_files(prs, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_header_bar(slide, "Skill Files  ·  Where Domain Knowledge Lives",
                   "Versioned markdown — not conversation history")

    add_text(slide, Inches(0.4), Inches(1.55), Inches(12.5), Inches(0.4),
             "Anatomy of a skill file", font_size=15, bold=True, color=NAVY)

    add_code_block(slide, Inches(0.4), Inches(2.0), Inches(7.5), Inches(4.0),
                   "---\n"
                   "name: treasury-iam\n"
                   "description: Treasury IAM specialist — FOE/BOE SoD,\n"
                   "             T_DEAL_*, T_TOE_HR, country variants\n"
                   "---\n\n"
                   "## When to use\n"
                   "- Validate FOE / BOE / MOE SoD compliance\n"
                   "- Plan catalog splits for forbidden combinations\n"
                   "- Audit hedge request management on T_TOE_HR\n\n"
                   "## Forbidden combinations\n"
                   "FOE: D3×01, D3×85, D2×AB\n"
                   "BOE: D2×01, D2×02, D2×16, D2×85, D2×KU, D2×VF\n\n"
                   "## Workflow\n"
                   "1. Look up the app in APS_IAM_W_APP\n"
                   "2. Find catalog via APS_IAM_W_BC_APP\n"
                   "3. Pull APPAUI then APPAUV\n"
                   "4. Cross-check forbidden matrix\n",
                   font_size=11)

    # Right — explainers
    add_text(slide, Inches(8.2), Inches(2.0), Inches(4.7), Inches(0.4),
             "Three sync targets", font_size=14, bold=True, color=NAVY)
    sync_items = [
        (".claude/skills/", "Source of truth — edit here."),
        ("skills/", "Mirror with frontmatter — for git diffs."),
        (".claude/commands/", "Slash command — frontmatter stripped."),
    ]
    add_bullets(slide, Inches(8.2), Inches(2.4), Inches(4.7), Inches(2.0),
                sync_items, font_size=11, bullet_color=ACCENT)

    add_text(slide, Inches(8.2), Inches(4.55), Inches(4.7), Inches(0.4),
             "Two activation modes", font_size=14, bold=True, color=NAVY)
    modes = [
        ("Explicit.", "Type /treasury-iam — full context loaded immediately."),
        ("Automatic.", "Keyword match (FOE, T_DEAL, hedge…) triggers selection."),
    ]
    add_bullets(slide, Inches(8.2), Inches(4.95), Inches(4.7), Inches(1.6),
                modes, font_size=11, bullet_color=GREEN)

    add_footer(slide, page_num, total)


def slide_treasury_sod(prs, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_header_bar(slide, "Treasury SoD  ·  Forbidden Matrices",
                   "Encoded in /treasury-iam — checked automatically on every analysis")

    # FOE table
    add_text(slide, Inches(0.4), Inches(1.55), Inches(6.0), Inches(0.4),
             "FOE catalog must NOT contain", font_size=15, bold=True, color=RED)
    foe_rows = [
        ["D3", "01", "Create settlement"],
        ["D3", "85", "Reverse settlement"],
        ["D2", "AB", "Settle contract"],
    ]
    add_table(slide, Inches(0.4), Inches(2.0), Inches(6.0), Inches(1.4),
              ["TRFCT", "ACTVT", "Meaning"], foe_rows,
              body_size=12, header_size=12, header_fill=RED)

    # BOE table
    add_text(slide, Inches(0.4), Inches(3.65), Inches(6.0), Inches(0.4),
             "BOE catalog must NOT contain", font_size=15, bold=True, color=RED)
    boe_rows = [
        ["D2", "01", "Create contract"],
        ["D2", "02", "Edit contract"],
        ["D2", "16", "Execute contract"],
        ["D2", "85", "Reverse contract"],
        ["D2", "KU", "Give Notice"],
        ["D2", "VF", "Expire contract"],
    ]
    add_table(slide, Inches(0.4), Inches(4.1), Inches(6.0), Inches(2.6),
              ["TRFCT", "ACTVT", "Meaning"], boe_rows,
              body_size=12, header_size=12, header_fill=RED)

    # Hedge SoD on right
    add_text(slide, Inches(6.85), Inches(1.55), Inches(6.0), Inches(0.4),
             "Hedge (T_TOE_HR)  ·  MOE forbidden", font_size=15, bold=True, color=RED)
    moe_rows = [
        ["A", "43", "Release manual designation"],
        ["A", "85", "Reverse manual designation"],
        ["D", "43", "Release dedesignation"],
        ["D", "85", "Reverse dedesignation"],
    ]
    add_table(slide, Inches(6.85), Inches(2.0), Inches(6.0), Inches(1.7),
              ["HREQ_CAT", "ACTVT", "Meaning"], moe_rows,
              body_size=12, header_size=12, header_fill=RED)

    add_text(slide, Inches(6.85), Inches(3.95), Inches(6.0), Inches(0.4),
             "Hedge (T_TOE_HR)  ·  Accountant forbidden", font_size=15, bold=True, color=RED)
    acc_rows = [
        ["A", "01", "Create manual designation"],
        ["A", "02", "Change manual designation"],
        ["A", "06", "Delete manual designation"],
        ["D", "01", "Create dedesignation"],
        ["D", "02", "Change dedesignation"],
        ["D", "06", "Delete dedesignation"],
    ]
    add_table(slide, Inches(6.85), Inches(4.4), Inches(6.0), Inches(2.6),
              ["HREQ_CAT", "ACTVT", "Meaning"], acc_rows,
              body_size=12, header_size=12, header_fill=RED)

    add_footer(slide, page_num, total)


def slide_treasury_brts(prs, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_header_bar(slide, "Treasury BRTs and Auth Objects",
                   "What roles look like, what objects they touch")

    # BRTs table
    add_text(slide, Inches(0.4), Inches(1.55), Inches(12.5), Inches(0.4),
             "Business Role Templates", font_size=15, bold=True, color=NAVY)
    brts = [
        ["SAP_BR_TREASURY_SPECIALIST_FOE", "Treasury Specialist · Front Office", "FOE",
         "Place orders, manage trading positions"],
        ["SAP_BR_TREASURY_SPECIALIST_BOE", "Treasury Specialist · Back Office", "BOE",
         "Settle contracts, post payments"],
        ["SAP_BR_TREASURY_SPECIALIST_MOE", "Treasury Specialist · Middle Office", "MOE",
         "Risk monitoring, hedge designation"],
        ["SAP_BR_TREASURY_ACCOUNTANT", "Treasury Accountant", "BOE",
         "Accounting close, dedesignation release"],
        ["SAP_BR_TREASURY_RISK_MANAGER", "Treasury Risk Manager", "Mixed",
         "Limit setting, exposure analysis"],
    ]
    add_table(slide, Inches(0.4), Inches(2.0), Inches(12.5), Inches(2.0),
              ["BRT ID", "Display Name", "Office", "Responsibility"], brts,
              body_size=10, header_size=11)

    # Country variants note
    add_text(slide, Inches(0.4), Inches(4.15), Inches(12.5), Inches(0.4),
             "Country variants  ·  _CN, _SG, _TH, _BR follow the same catalog pattern as the base BRT.",
             font_size=11, color=GREY_MED)

    # Core auth objects table
    add_text(slide, Inches(0.4), Inches(4.55), Inches(12.5), Inches(0.4),
             "Core auth objects", font_size=15, bold=True, color=NAVY)
    auth_rows = [
        ["T_DEAL_PD", "Company / product / transaction type", "All transactions",
         "TRFCT, ACTVT, BUKRS"],
        ["T_DEAL_PF", "Company / portfolio", "All transactions",
         "TRFCT, ACTVT, BUKRS, RPORTB"],
        ["T_DEAL_DP", "Company / securities account", "Securities only",
         "TRFCT, ACTVT, BUKRS, RDEPO"],
        ["T_DEAL_AG", "Company / authorization group", "Unused in cloud",
         "TRFCT, ACTVT, BUKRS"],
        ["T_TOE_HR", "Hedge request management", "Hedge designation flow",
         "HREQ_CAT, ACTVT"],
    ]
    add_table(slide, Inches(0.4), Inches(5.0), Inches(12.5), Inches(2.0),
              ["Auth Object", "Scope", "Applies to", "Key fields"], auth_rows,
              body_size=10, header_size=11)
    add_footer(slide, page_num, total)


def slide_cash_detail(prs, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_header_bar(slide, "Cash Management  ·  Auth Objects & Activity Matrix",
                   "Bank account governance and the four-eyes principle")

    # Auth objects table
    add_text(slide, Inches(0.4), Inches(1.55), Inches(12.5), Inches(0.4),
             "Authorization objects", font_size=15, bold=True, color=NAVY)
    objs = [
        ["F_CLM_BAM", "Bank account master data",
         "01=Create, 02=Change, 03=Display, 06=Delete, 63=Transport"],
        ["F_CLM_BAOR", "Bank account opening request",
         "03=Display, 31=Approve"],
        ["F_CLM_BKCR", "Bank account change request",
         "01=Create, 02=Change, 03=Display"],
        ["F_CLM_BAI", "Bank account interest records",
         "03=Display, 06=Delete"],
        ["F_CLM_BAIC", "Bank account interest conditions",
         "01=Create, 02=Change, 03=Display, 06=Delete, F4=Value Help"],
        ["F_CLM_BAH2", "Bank account hierarchy",
         "01=Create, 02=Change, 03=Display, 06=Delete"],
    ]
    add_table(slide, Inches(0.4), Inches(2.0), Inches(12.5), Inches(2.4),
              ["Auth Object", "Description", "Key ACTVT values"], objs,
              body_size=10, header_size=11)

    # Activity matrix
    add_text(slide, Inches(0.4), Inches(4.55), Inches(12.5), Inches(0.4),
             "Activity matrix per key Cash Management app",
             font_size=15, bold=True, color=NAVY)
    matrix = [
        ["F1366_TRAN", "Manage Bank Accounts", "01, 02, 03, 06", "—", "—"],
        ["F1366A_TRAN", "Display Bank Accounts", "03", "—", "—"],
        ["F5861_TRAN", "Submit Applications", "—", "(submit flow)", "—"],
        ["F5859_TRAN", "Approve Applications", "—", "03, 31", "—"],
        ["F5860_TRAN", "Applications Overview", "—", "03, 31", "—"],
        ["F6264_TRAN", "Approve Changes", "03, 63", "—", "01, 02, 03"],
        ["F9015_TRAN", "Monitor Interest", "03", "—", "—"],
        ["F9017_TRAN", "Manage Interest Conditions", "03", "—", "—"],
    ]
    add_table(slide, Inches(0.4), Inches(5.0), Inches(12.5), Inches(2.0),
              ["App", "Description", "F_CLM_BAM", "F_CLM_BAOR", "F_CLM_BKCR"], matrix,
              body_size=10, header_size=11)
    add_footer(slide, page_num, total)


def slide_finance_detail(prs, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_header_bar(slide, "Finance  ·  AP / AR / GL / BA",
                   "Auth objects and SoD touchpoints across the FI flows")

    headers = ["Sub-domain", "Package", "Core auth objects", "SoD focus"]
    rows = [
        ["Accounts Payable (AP)", "CLOUD_FI_AP_IAM",
         "F_BKPF_BUK / BLA / KOA, F_LFA1_BUK",
         "Invoice creation ≠ payment approval (four-eyes)"],
        ["Accounts Receivable (AR)", "CLOUD_FI_AR_IAM",
         "F_BKPF_BUK / BLA / KOA, F_KNA1_BUK",
         "Dunning ≠ incoming payment posting"],
        ["General Ledger (GL)", "CLOUD_FI_GL_IAM",
         "F_BKPF_BUK / BLA / KOA",
         "Posting ≠ verification ≠ park-and-release"],
        ["Bank Accounting (BA)", "CLOUD_FI_BA_IAM",
         "F_BKPF_BUK + bank statement objects",
         "Statement upload ≠ posting ≠ reconciliation"],
    ]
    add_table(slide, Inches(0.4), Inches(1.55), Inches(12.5), Inches(2.4),
              headers, rows, body_size=10, header_size=11)

    # AP payment flow detail
    add_text(slide, Inches(0.4), Inches(4.1), Inches(12.5), Inches(0.4),
             "Example  ·  AP payment-run SoD flow",
             font_size=15, bold=True, color=NAVY)

    flow = [
        ("F0770_TRAN", "Schedule Payment Proposal", "Propose"),
        ("F0771_TRAN", "Revise Payment Proposal", "Revise"),
        ("F0673A_TRAN", "Approve Payment Run", "Approve"),
    ]
    box_w_in = 3.7
    gap_in = 0.4
    start_in = (13.333 - (box_w_in * 3 + gap_in * 2)) / 2
    top = Inches(4.6)

    for i, (app, desc, role) in enumerate(flow):
        left = Inches(start_in + (box_w_in + gap_in) * i)
        add_rect(slide, left, top, Inches(box_w_in), Inches(1.6),
                 fill=WHITE, line=GREY_LIGHT)
        add_rect(slide, left, top, Inches(box_w_in), Inches(0.45), fill=ACCENT)
        add_text(slide, left + Inches(0.15), top + Inches(0.08),
                 Inches(box_w_in - 0.3), Inches(0.32),
                 role, font_size=12, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)
        add_text(slide, left + Inches(0.15), top + Inches(0.55),
                 Inches(box_w_in - 0.3), Inches(0.4),
                 app, font_size=14, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER, font_name="Consolas")
        add_text(slide, left + Inches(0.15), top + Inches(0.95),
                 Inches(box_w_in - 0.3), Inches(0.5),
                 desc, font_size=11, color=GREY_DARK,
                 align=PP_ALIGN.CENTER)

        if i < 2:
            arrow_left = Inches(start_in + (box_w_in + gap_in) * i + box_w_in)
            arr = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                arrow_left + Inches(0.05),
                top + Inches(0.6),
                Inches(gap_in - 0.1),
                Inches(0.4),
            )
            arr.shadow.inherit = False
            arr.fill.solid()
            arr.fill.fore_color.rgb = GREY_MED
            arr.line.fill.background()

    add_text(slide, Inches(0.4), Inches(6.4), Inches(12.5), Inches(0.4),
             "All three apps must NOT be held by the same user without compensating control.",
             font_size=11, color=RED, align=PP_ALIGN.CENTER, bold=True)
    add_footer(slide, page_num, total)


def slide_query_walkthrough(prs, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_header_bar(slide, "Query Walk-through  ·  FTRCAI02_B_TRAN",
                   "The actual SQL Claude runs to validate one Treasury app")

    # Step 1
    add_text(slide, Inches(0.4), Inches(1.5), Inches(12.5), Inches(0.4),
             "Step 1  ·  Find catalog and BRT for the app",
             font_size=12, bold=True, color=ACCENT)
    add_code_block(slide, Inches(0.4), Inches(1.85), Inches(12.5), Inches(1.0),
                   "SELECT BU_CATALOG_ID, APP_ID FROM APS_IAM_W_BC_APP\n"
                   "  WHERE APP_ID = 'FTRCAI02_B_TRAN'   -- rows: 5\n",
                   font_size=11)

    # Step 2
    add_text(slide, Inches(0.4), Inches(2.95), Inches(12.5), Inches(0.4),
             "Step 2  ·  Pull active auth-object instances",
             font_size=12, bold=True, color=ACCENT)
    add_code_block(slide, Inches(0.4), Inches(3.3), Inches(12.5), Inches(1.0),
                   "SELECT UUID, AUTH_OBJECT FROM APS_IAM_W_APPAUI\n"
                   "  WHERE APP_ID = 'FTRCAI02_B_TRAN' AND INACTIVE = ' '   -- rows: 30\n",
                   font_size=11)

    # Step 3
    add_text(slide, Inches(0.4), Inches(4.4), Inches(12.5), Inches(0.4),
             "Step 3  ·  Pull TRFCT and ACTVT field values for those instances",
             font_size=12, bold=True, color=ACCENT)
    add_code_block(slide, Inches(0.4), Inches(4.75), Inches(12.5), Inches(1.4),
                   "SELECT PARENT_ID, FIELD, LOW_VALUE, HIGH_VALUE\n"
                   "  FROM APS_IAM_W_APPAUV\n"
                   "  WHERE APP_ID = 'FTRCAI02_B_TRAN'\n"
                   "    AND FIELD IN ('TRFCT', 'ACTVT')   -- rows: 50\n",
                   font_size=11)

    # Step 4 - verdict
    add_text(slide, Inches(0.4), Inches(6.25), Inches(12.5), Inches(0.4),
             "Step 4  ·  Cross-check held (TRFCT, ACTVT) pairs against the BOE forbidden matrix in memory.  ✓ COMPLIANT",
             font_size=12, bold=True, color=GREEN)
    add_footer(slide, page_num, total)


def slide_execute_report(prs, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_header_bar(slide, "/execute  ·  Structured Output Report",
                   "Every autonomous run produces the same five sections")

    sections = [
        ("Goal", ACCENT, "Restate the original objective verbatim — what success looks like."),
        ("Summary", ACCENT2, "One-paragraph executive summary  ·  scope, what was checked, headline verdict."),
        ("Findings", NAVY, "Numbered list of facts: app/auth-object/value combinations actually present in ER6."),
        ("Violations", RED, "Subset of findings that breach a rule  ·  cite the rule and the specific value."),
        ("Recommendations", GREEN, "Concrete next actions  ·  catalog split, value removal, BRT realignment."),
    ]
    top = 1.65
    for i, (label, color, body) in enumerate(sections):
        y = top + i * 1.05
        add_rect(slide, Inches(0.4), Inches(y), Inches(2.2), Inches(0.85), fill=color)
        add_text(slide, Inches(0.4), Inches(y + 0.1), Inches(2.2), Inches(0.65),
                 label, font_size=18, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(slide, Inches(2.7), Inches(y), Inches(10.2), Inches(0.85),
                 fill=WHITE, line=GREY_LIGHT)
        add_text(slide, Inches(2.85), Inches(y + 0.1), Inches(9.9), Inches(0.65),
                 body, font_size=12, color=GREY_DARK, anchor=MSO_ANCHOR.MIDDLE)

    add_text(slide, Inches(0.4), Inches(6.95), Inches(12.5), Inches(0.4),
             "Every claim is backed by the queries logged to .session-log.md  →  fully reproducible audit trail.",
             font_size=11, color=GREY_MED, align=PP_ALIGN.CENTER)
    add_footer(slide, page_num, total)


def slide_memo_system(prs, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_header_bar(slide, "/memo  ·  Persistent Investigation Memory",
                   "Findings survive across sessions and analysts")

    # Sub-commands
    add_text(slide, Inches(0.4), Inches(1.55), Inches(6.0), Inches(0.4),
             "Sub-commands", font_size=14, bold=True, color=NAVY)
    cmd_rows = [
        ["/memo save [topic]", "Save current session findings"],
        ["/memo update [topic]", "Append new findings to existing memo"],
        ["/memo load [topic]", "Reload a memo and offer to resume"],
        ["/memo list", "Show all memos from INDEX.md"],
        ["/memo clear <topic>", "Delete a memo (asks for confirmation)"],
    ]
    add_table(slide, Inches(0.4), Inches(2.0), Inches(6.0), Inches(2.6),
              ["Command", "Action"], cmd_rows, body_size=11, header_size=12)

    # Required sections
    add_text(slide, Inches(0.4), Inches(4.8), Inches(6.0), Inches(0.4),
             "Every memo  ·  4 required sections", font_size=14, bold=True, color=NAVY)
    sec_items = [
        ("## Findings", "Facts established this session."),
        ("## Decisions", "Choices made and the why."),
        ("## Work in Progress", "Open threads, next steps."),
        ("## Known Good Baselines", "What is verified-correct."),
    ]
    add_bullets(slide, Inches(0.4), Inches(5.25), Inches(6.0), Inches(1.7),
                sec_items, font_size=11, bullet_color=ACCENT)

    # Right panel — auto-resume
    add_text(slide, Inches(6.85), Inches(1.55), Inches(6.0), Inches(0.4),
             "Auto-surfacing on session start",
             font_size=14, bold=True, color=NAVY)
    add_rect(slide, Inches(6.85), Inches(2.0), Inches(6.0), Inches(2.0),
             fill=WHITE, line=GREY_LIGHT)
    add_text(slide, Inches(7.0), Inches(2.15), Inches(5.7), Inches(1.7),
             "Found in-progress memo:\n"
             "    trm-catalog-split\n"
             "Last updated 2026-05-10.\n"
             "Load it to resume?",
             font_size=13, color=GREY_DARK)

    add_text(slide, Inches(6.85), Inches(4.2), Inches(6.0), Inches(0.4),
             "End-to-end workflow", font_size=14, bold=True, color=NAVY)
    flow_items = [
        "/goal Validate FOE/BOE for CLOUD_FI_TR_IAM",
        "    ↓",
        "/execute  →  produces report",
        "    ↓",
        "/memo save trm-cloud-foe-boe",
        "    ↓ (next session)",
        "/memo load trm-cloud-foe-boe  →  resume",
    ]
    box = slide.shapes.add_textbox(Inches(6.85), Inches(4.65), Inches(6.0), Inches(2.4))
    tf = box.text_frame
    tf.word_wrap = True
    for idx, item in enumerate(flow_items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.line_spacing = 1.15
        run = p.add_run()
        run.text = item
        run.font.name = "Consolas"
        run.font.size = Pt(11)
        run.font.color.rgb = GREY_DARK
    add_footer(slide, page_num, total)


def slide_hooks_detail(prs, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_header_bar(slide, "Quality Hooks  ·  Deep Dive",
                   "Three bash hooks that run silently every session")

    # Hook 1
    add_card(slide, Inches(0.4), Inches(1.6), Inches(4.1), Inches(5.0),
             title="validate-memo.sh",
             body=(
                 "Phase: PreToolUse / Write | Edit\n"
                 "Match: .claude/memo/*.md\n\n"
                 "What it does\n"
                 "  •  Reads the proposed file content\n"
                 "  •  Greps for the 4 required sections\n"
                 "  •  Exits non-zero if any are missing\n"
                 "  •  Reports which sections were missing\n"
                 "  •  Edit short-circuits — partial replace\n"
                 "    is validated on the next full Write\n\n"
                 "Why it matters\n"
                 "  Prevents shallow / abandoned memos\n"
                 "  Enforces uniform structure across\n"
                 "  every investigation."
             ),
             accent=ACCENT, body_size=11)

    # Hook 2
    add_card(slide, Inches(4.65), Inches(1.6), Inches(4.1), Inches(5.0),
             title="sync-skills.sh",
             body=(
                 "Phase: PostToolUse / Write | Edit\n"
                 "Match: .claude/skills/*.md\n\n"
                 "What it does\n"
                 "  •  Copies edited skill to skills/\n"
                 "    (with frontmatter)\n"
                 "  •  Copies it to .claude/commands/\n"
                 "    (frontmatter stripped)\n"
                 "  •  Three locations stay identical\n"
                 "  •  Both Write and Edit are matched\n"
                 "    so partial edits stay in sync\n\n"
                 "Why it matters\n"
                 "  Skills are the source of domain truth.\n"
                 "  Drift = silent inconsistency."
             ),
             accent=ACCENT2, body_size=11)

    # Hook 3
    add_card(slide, Inches(8.9), Inches(1.6), Inches(4.0), Inches(5.0),
             title="log-query.sh",
             body=(
                 "Phase: PostToolUse / query_sql\n"
                 "Match: mcp__er6__query_sql\n\n"
                 "What it does\n"
                 "  •  Captures every SQL run\n"
                 "  •  Appends to .session-log.md async\n"
                 "  •  Records timestamp + row limit\n\n"
                 "Why it matters\n"
                 "  Reproducibility  ·  every claim ties\n"
                 "  back to the exact query that\n"
                 "  produced it.\n"
                 "  Audit trail — who asked what when."
             ),
             accent=GREEN, body_size=11)

    add_footer(slide, page_num, total)


def slide_webui_detail(prs, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_header_bar(slide, "Web UI  ·  Stack and Features",
                   "Browser chat for the rest of the team")

    # Stack table
    add_text(slide, Inches(0.4), Inches(1.55), Inches(6.0), Inches(0.4),
             "Stack", font_size=14, bold=True, color=NAVY)
    stack = [
        ["Web framework", "FastAPI + Uvicorn"],
        ["LLM streaming", "Anthropic SDK · SSE · Claude Opus"],
        ["MCP tools", "Subprocess over stdio — er6 (6 tools) + sap-wiki (4 tools)"],
        ["Auth", "OIDC (Authlib) · signed session cookie"],
        ["Templates", "Jinja2  ·  ui/templates/"],
        ["Static assets", "ui/static/  (CSS + JS, no bundler)"],
        ["Markdown", "marked + DOMPurify  ·  vendored locally"],
        ["Theming", "CSS variables · data-theme · Catppuccin"],
    ]
    add_table(slide, Inches(0.4), Inches(2.0), Inches(6.0), Inches(4.5),
              ["Layer", "Technology"], stack, body_size=11, header_size=12)

    # Features
    add_text(slide, Inches(6.85), Inches(1.55), Inches(6.0), Inches(0.4),
             "Features", font_size=14, bold=True, color=NAVY)
    features = [
        ("Streaming chat.", "Tokens render as they arrive; tool calls expand inline."),
        ("Theme toggle.", "Light = Catppuccin Latte · Dark = Mocha · persisted in localStorage."),
        ("No-flash theme.", "Inline <head> script applies saved theme before stylesheet."),
        ("Prompt template library.", "Categorised welcome screen with 30+ examples."),
        ("Session history.", "Multiple sessions per user with rename/delete."),
        ("OIDC + dev mode.", "Real OIDC in prod  ·  auto-login as 'dev' for local."),
        ("Tool call transparency.", "Every MCP call is shown — table, params, result rows."),
        ("Wiki access.", "/iam-wiki and SAP Confluence available in both CLI and Web UI via sap-wiki MCP."),
        ("Connection status.", "Header shows ER6 and Wiki status independently (green = connected)."),
        ("Vendored libs.", "marked + DOMPurify shipped locally — zero CDN dependency."),
    ]
    add_bullets(slide, Inches(6.85), Inches(2.0), Inches(6.0), Inches(4.8),
                features, font_size=10, bullet_color=ACCENT)
    add_footer(slide, page_num, total)


def slide_prompt_templates(prs, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_header_bar(slide, "Prompt Template Library",
                   "Categorised welcome-screen examples to bootstrap analysts")

    headers = ["Category", "Template Examples"]
    rows = [
        ["Getting Started",
         "What can I ask?  ·  Glossary BRT/BC/App  ·  What is SoD?  ·  How do roles work?"],
        ["General",
         "App → Catalog mapping  ·  Restriction type coverage  ·  BRT catalog tree  ·  Auth object usage"],
        ["Finance — AP",
         "App auth analysis  ·  AP payment flow SoD  ·  Vendor master access"],
        ["Finance — AR",
         "AR clearing & incoming payment  ·  AR dunning SoD check  ·  Customer master access"],
        ["Finance — GL",
         "GL posting controls  ·  Activity set check  ·  Document type access"],
        ["Finance — BA",
         "Bank statement health check  ·  Reconciliation auth review"],
        ["Treasury and Risk",
         "FOE/BOE SoD validation  ·  Catalog split analysis  ·  Hedge request SoD  ·  BRT footprint"],
        ["Cash Management",
         "Activity set completeness  ·  Submit/Approve SoD  ·  Four-eyes catalog check  ·  IAM health"],
    ]
    add_table(slide, Inches(0.4), Inches(1.55), Inches(12.5), Inches(5.4),
              headers, rows, body_size=11, header_size=12)
    add_footer(slide, page_num, total)


def slide_performance(prs, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_header_bar(slide, "Performance  ·  Measured Against Live ER6",
                   "Validated 2026-05-15 · 20 tests · 100% pass rate")

    # Tool timings
    add_text(slide, Inches(0.4), Inches(1.55), Inches(6.0), Inches(0.4),
             "Round-trip latency by tool", font_size=14, bold=True, color=NAVY)
    tools = [
        ["query_sql", "17", "10,106 ms", "8,728", "12,603"],
        ["read_table_def", "2", "8,922 ms", "8,650", "9,193"],
        ["list_package", "1", "17,322 ms", "—", "—"],
    ]
    add_table(slide, Inches(0.4), Inches(2.0), Inches(6.0), Inches(2.0),
              ["Tool", "Calls", "Mean", "Min", "Max"], tools,
              body_size=11, header_size=11)

    # Selected operation timings
    add_text(slide, Inches(6.85), Inches(1.55), Inches(6.0), Inches(0.4),
             "Selected operations", font_size=14, bold=True, color=NAVY)
    ops = [
        ["TADIR · 20 SIA6 objects", "12,603 ms"],
        ["list_package · 700+ objects", "17,322 ms"],
        ["APS_IAM_W_APP · 1 row", "9,592 ms"],
        ["APS_IAM_W_APP · 20 LIKE", "10,732 ms"],
        ["APPAUI · 5–15 rows", "9,129–10,082 ms"],
        ["APPAUV · 21–50 rows", "9,567–12,001 ms"],
        ["BC_APP · 1 row", "9,343–11,053 ms"],
        ["BRTBUC · 29–39 rows", "9,131–11,223 ms"],
        ["read_table_def (DDL only)", "8,650–9,193 ms"],
    ]
    add_table(slide, Inches(6.85), Inches(2.0), Inches(6.0), Inches(4.4),
              ["Operation", "Duration"], ops, body_size=10, header_size=11)

    # Notes
    add_text(slide, Inches(0.4), Inches(4.2), Inches(6.0), Inches(0.4),
             "What this means", font_size=14, bold=True, color=NAVY)
    notes = [
        ("Typical investigation:", "5–8 queries  ≈  50–90 s ER6 round-trip."),
        ("list_package is slowest.", "Avoid on 700+ object packages unless needed."),
        ("read_table_def is fastest.", "DDL only, no table scan."),
        ("Row count has modest impact.", "50 rows ≈ 2.5 s slower than 1 row."),
        ("Server-side cache:", "Repeat queries run ~19% faster (11.2s → 9.1s)."),
    ]
    add_bullets(slide, Inches(0.4), Inches(4.6), Inches(6.0), Inches(2.4),
                notes, font_size=11, bullet_color=ACCENT)
    add_footer(slide, page_num, total)


def slide_security(prs, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_header_bar(slide, "Security & Compliance",
                   "Read-only by construction · auditable by design")

    # Pillars
    pillars = [
        ("Read-only access", ACCENT,
         "All ER6 access goes through the ANZEIGER user — display privileges only. "
         "No SAP authorization to write, delete, or release transports. "
         "Tool can run during business hours without change-control concerns."),
        ("Encrypted transport", ACCENT2,
         "SSL between MCP server and ER6 (sapcli config). "
         "Anthropic traffic over TLS. OIDC sessions over HTTPS in production. "
         "Session cookies signed with SESSION_SECRET."),
        ("Authentication", GREEN,
         "Web UI uses OIDC (Authlib). Each user authenticates with their corporate IdP; "
         "session cookie binds queries to a user identity. Dev mode (auto-login as 'dev') "
         "is explicitly opt-in via OIDC_CLIENT_ID placeholder."),
        ("Audit trail", NAVY,
         "Every ER6 SQL query is appended to .session-log.md with a timestamp by the "
         "log-query hook. Memo files preserve findings per investigation. "
         "All answers cite the queries that produced them."),
        ("No data egress", RED,
         "Only the LLM provider receives the question, the queries Claude composes, "
         "and the row data needed to answer. No bulk dumps; row limits are explicit. "
         "All data stays on ER6 except what Claude needs to reason about a question."),
    ]
    top = 1.55
    h = 1.05
    for i, (title, color, body) in enumerate(pillars):
        y = top + i * h
        add_rect(slide, Inches(0.4), Inches(y), Inches(2.6), Inches(h - 0.05), fill=color)
        add_text(slide, Inches(0.4), Inches(y + 0.1), Inches(2.6), Inches(h - 0.25),
                 title, font_size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(slide, Inches(3.1), Inches(y), Inches(9.8), Inches(h - 0.05),
                 fill=WHITE, line=GREY_LIGHT)
        add_text(slide, Inches(3.25), Inches(y + 0.08), Inches(9.55), Inches(h - 0.2),
                 body, font_size=11, color=GREY_DARK, anchor=MSO_ANCHOR.MIDDLE)
    add_footer(slide, page_num, total)


def slide_directory(prs, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_header_bar(slide, "Project Structure",
                   "Where everything lives in the repository")

    add_code_block(slide, Inches(0.4), Inches(1.55), Inches(12.5), Inches(5.4),
                   "iam-assistant/\n"
                   "├── CLAUDE.md                     # Claude Code instructions  ·  data dictionary\n"
                   "├── README.md\n"
                   "├── pyproject.toml\n"
                   "├── .sapcli.env                   # ER6 credentials  (not committed; fallback only)\n"
                   "├── .env                          # Web UI settings  (not committed)\n"
                   "├── .mcp.json                     # MCP server registration\n"
                   "├── app/                          # Web UI application\n"
                   "│   ├── main.py                   # FastAPI app  ·  routes  ·  MCP lifespan\n"
                   "│   ├── auth.py                   # OIDC routes  ·  session dependency\n"
                   "│   ├── chat.py                   # Anthropic streaming  ·  tool execution\n"
                   "│   ├── config.py                 # Pydantic settings\n"
                   "│   └── mcp_client.py             # MCP subprocess client (stdio)\n"
                   "├── ui/\n"
                   "│   ├── templates/index.html\n"
                   "│   └── static/  app.js · style.css · prompt-templates.js · vendor/\n"
                   "├── mcp-server/\n"
                   "│   └── er6_mcp_server.py         # MCP server exposing six ER6 tools\n"
                   "├── skills/  (mirror)             # treasury-iam · cash-iam · fin-iam · goal · execute · memo\n"
                   "├── tests/                        # pytest suite (chat / auth / mcp / config / main)\n"
                   "├── docs/                         # this deck lives here\n"
                   "└── .claude/\n"
                   "    ├── hooks/  validate-memo · sync-skills · log-query\n"
                   "    ├── memo/   INDEX.md · .session-log.md · per-investigation memos\n"
                   "    ├── skills/    (source of truth)\n"
                   "    ├── commands/  (slash-command mirror)\n"
                   "    └── settings.json\n",
                   font_size=10)
    add_footer(slide, page_num, total)


def slide_lessons(prs, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_header_bar(slide, "Lessons Learned",
                   "What we discovered building this")

    left = [
        ("Skills > prompts.", "Versioned skill files keep domain knowledge out of conversation history."),
        ("Encode the gotchas.", "Every undocumented quirk (BUC unqueryable, GUID app IDs) goes into the skill."),
        ("Validation matters.", "20 live-data tests caught query patterns that looked right but failed at runtime."),
        ("Hooks are leverage.", "Three small bash scripts enforce more discipline than any policy doc."),
    ]
    add_text(slide, Inches(0.4), Inches(1.55), Inches(6.1), Inches(0.4),
             "Engineering takeaways", font_size=14, bold=True, color=NAVY)
    add_bullets(slide, Inches(0.4), Inches(2.0), Inches(6.1), Inches(4.5),
                left, font_size=12, bullet_color=ACCENT)

    right = [
        ("Read-only is liberating.", "No change control, no risk — analysts use it freely during operations."),
        ("Time-to-answer changes behavior.", "When SoD checks are 60 s instead of hours, people actually run them."),
        ("Audit trail builds trust.", "Every claim shows its query — reviewers verify, not just believe."),
        ("Persistent memo wins.", "Multi-day investigations resume cleanly across sessions and analysts."),
    ]
    add_text(slide, Inches(6.75), Inches(1.55), Inches(6.1), Inches(0.4),
             "Adoption takeaways", font_size=14, bold=True, color=NAVY)
    add_bullets(slide, Inches(6.75), Inches(2.0), Inches(6.1), Inches(4.5),
                right, font_size=12, bullet_color=GREEN)

    add_footer(slide, page_num, total)


def slide_faq(prs, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_LIGHT)
    add_header_bar(slide, "FAQ",
                   "Common questions from new users")

    qa = [
        ("Q. Can the assistant modify ER6 data?",
         "A. No. The ANZEIGER user has display-only privileges. Even if asked, no write is possible."),
        ("Q. Is my prompt sent to Anthropic?",
         "A. Routed through Hyperspace at localhost:6655. Configure ANTHROPIC_BASE_URL in .env."),
        ("Q. What if a query fails or times out?",
         "A. Claude reports the failure, retries with a smaller scope, or falls back to read_table_def to verify column names."),
        ("Q. How do I add a new domain (e.g. Procurement)?",
         "A. Author a new skill file in .claude/skills/. The sync hook propagates it to skills/ and .claude/commands/."),
        ("Q. Can two analysts use the Web UI simultaneously?",
         "A. Yes — each OIDC session is independent; the MCP server multiplexes calls."),
        ("Q. Does it work offline?",
         "A. No. ER6 connectivity is required (and an Anthropic-compatible endpoint)."),
        ("Q. What happens if a skill conflicts with CLAUDE.md?",
         "A. Explicit user instructions and CLAUDE.md take precedence over skill content."),
        ("Q. Can I run only the CLI without the Web UI?",
         "A. Yes — the CLI uses only the MCP server. The Web UI is a separate optional FastAPI app."),
    ]
    box = slide.shapes.add_textbox(Inches(0.4), Inches(1.55), Inches(12.5), Inches(5.4))
    tf = box.text_frame
    tf.word_wrap = True
    for idx, (q, a) in enumerate(qa):
        if idx > 0:
            p = tf.add_paragraph()
            p.space_before = Pt(6)
        p_q = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p_q.line_spacing = 1.15
        r = p_q.add_run()
        r.text = q
        r.font.name = "Calibri"
        r.font.bold = True
        r.font.size = Pt(12)
        r.font.color.rgb = NAVY

        p_a = tf.add_paragraph()
        p_a.line_spacing = 1.15
        p_a.space_after = Pt(2)
        r = p_a.add_run()
        r.text = a
        r.font.name = "Calibri"
        r.font.size = Pt(11)
        r.font.color.rgb = GREY_DARK
    add_footer(slide, page_num, total)


# ---------- Main ----------

def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        # --- Opening ---
        slide_title,
        slide_agenda,
        slide_glossary,

        # --- Part 1: Why ---
        lambda p, n, t: slide_section(p, n, t, "PART 1", "The Problem and the Vision"),
        slide_problem,
        slide_what_is,

        # --- Part 2: How it works ---
        lambda p, n, t: slide_section(p, n, t, "PART 2", "How It Works", color=ACCENT2),
        slide_architecture,
        slide_mcp_tools,
        slide_data_dict_overview,
        slide_table_relationships,
        slide_cds_view_layer,
        slide_skill_files,

        # --- Part 3: Domain skills ---
        lambda p, n, t: slide_section(p, n, t, "PART 3", "Domain Skills in Depth", color=GREEN),
        slide_skills,
        slide_treasury_sod,
        slide_treasury_brts,
        slide_cash_detail,
        slide_finance_detail,

        # --- Part 4: Workflow ---
        lambda p, n, t: slide_section(p, n, t, "PART 4", "Workflow Tools", color=NAVY),
        slide_workflow,
        slide_execute_report,
        slide_memo_system,
        slide_hooks_detail,

        # --- Part 5: Examples & operations ---
        lambda p, n, t: slide_section(p, n, t, "PART 5", "Examples and Operations", color=ACCENT),
        slide_example,
        slide_query_walkthrough,
        slide_use_cases,
        slide_webui_detail,
        slide_prompt_templates,
        slide_performance,
        slide_security,

        # --- Part 6: Wrap-up ---
        lambda p, n, t: slide_section(p, n, t, "PART 6", "Wrap-up", color=ACCENT2),
        slide_benefits,
        slide_lessons,
        slide_quickstart,
        slide_directory,
        slide_faq,
        slide_roadmap,
        slide_thanks,
    ]
    total = len(builders)
    for idx, builder in enumerate(builders, start=1):
        builder(prs, idx, total)

    out_dir = Path(__file__).resolve().parent.parent / "docs"
    out_dir.mkdir(exist_ok=True)
    out_path = out_dir / "IAM_Assistant_Overview.pptx"
    prs.save(out_path)
    print(f"Wrote {out_path}  ({total} slides)")


if __name__ == "__main__":
    build()
